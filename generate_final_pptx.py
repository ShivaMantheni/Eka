import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

def create_presentation():
    prs = Presentation()
    
    # Title Slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Eka Automation Application"
    subtitle.text = "User Guide & Demo Walkthrough"

    # Overview Slide
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    shapes.title.text = "Overview"
    body_shape = shapes.placeholders[1]
    tf = body_shape.text_frame
    tf.text = "The Eka Automation Application is designed to significantly reduce repetitive engineering workload through intelligent, centralized automation."
    
    p = tf.add_paragraph()
    p.text = "Daily Image Loading - Automated, zero-touch batch processing of VS testing images replacing manual CLI operations."
    
    p = tf.add_paragraph()
    p.text = "Batch Runs - Streamlined execution of multiple test configurations defined via YAML testbeds."
    
    p = tf.add_paragraph()
    p.text = "Advanced Regression Testing:"
    
    p2 = tf.add_paragraph()
    p2.text = "Parallel Execution: Run multiple SpyTest scripts simultaneously against discrete topologies."
    p2.level = 1
    
    p3 = tf.add_paragraph()
    p3.text = "Optimized DUT Usage: The system automatically manages Device Under Test (DUT) allocation and prevents resource conflicts."
    p3.level = 1
    
    p4 = tf.add_paragraph()
    p4.text = "Centralized Logging: Real-time streaming and storage of raw execution logs for rapid debugging and RCA."
    p4.level = 1


    slides_data = [
        {
            "title": "Dashboard",
            "text": ["Total Devices Dashboard: The central control plane landing page providing an immediate pulse check on your lab infrastructure."],
            "image": "screenshot_dashboard.png"
        },
        {
            "title": "Dashboard: Core Metrics",
            "text": [
                "4 Core Metrics:",
                "Total Devices: Total sum of registered IPs/Hostnames in the application database.",
                "Online: Result of the live background heartbeat ping checking node reachability.",
                "Scripts: Total test case scripts discovered from the connected Git Spytest repositories.",
                "Executions: Lifelong test runs triggered through the platform."
            ],
            "image": "spotlight_dashboard_1.png"
        },
        {
            "title": "Devices",
            "text": ["Manage the entire testing infrastructure inventory from this single pane of glass."],
            "image": "screenshot_devices.png"
        },
        {
            "title": "Devices: Provisioning",
            "text": [
                "Adding a New Device: This top section handles the secure ingestion of new testbed members into the database.",
                "Name / IP / Port: System requires unique Device Name string. IP Validated. Port field auto-defaults to 22.",
                "Credentials: Standard login parameters. Passwords encrypted at rest.",
                "Type Selection Dropdown: VM (Host), DUT, or Switch/Router."
            ],
            "image": "spotlight_devices_1.png"
        },
        {
            "title": "Devices: Management",
            "text": [
                "Inventory & Operations: The lower grid provides a comprehensive, filtered layout of the entire mapped lab topology.",
                "Live Status Checks: The Status column represents the last known heartbeat. Wi-Fi icon triggers on-demand validation.",
                "Decommissioning: Trash bin instantly purges records."
            ],
            "image": "spotlight_devices_2.png"
        },
        {
            "title": "Execution",
            "text": ["The core engine room. Configure, map, and dispatch complex testing pipelines."],
            "image": "screenshot_execute.png"
        },
        {
            "title": "Execution: Git Sync",
            "text": [
                "Host Server Selection: Bind the script to a physical host before proceeding.",
                "Git Orchestration Layer: Provide Git Clone URL, Access Token, and Branch. Trigger backend worker to SSH, clone, and index scripts."
            ],
            "image": "spotlight_execute_1.png"
        },
        {
            "title": "Execution: Topologies",
            "text": [
                "Dynamic Mapping: Selected DUTs spawn as draggable interface nodes on the canvas.",
                "Cable Mode Routing: Enter linking state, drag from source port to destination port to draw a physical link line, resetting as needed."
            ],
            "image": "spotlight_execute_2.png"
        },
        {
            "title": "Execution: Target Setup",
            "text": [
                "Script Selection (Center Panel): Select Category to auto-fill Test Scripts. Check multiple scripts to batch regression test queue.",
                "Testbed Definition (Right Panel): Select the YAML Testbed file instructing the SpyTest engine on layout."
            ],
            "image": "spotlight_execute_3.png"
        },
        {
            "title": "VS Manager",
            "text": ["Automates the painful process of updating firmware images across massive virtual simulation banks."],
            "image": "screenshot_vs.png"
        },
        {
            "title": "VS Manager: Update Workflow",
            "text": [
                "Discovery & Targeting: Queries the host server and spins up a list of all active Virtual Machines detected. Select VMs.",
                "Update Lifecycle Workflow: Enter the absolute file path to the fresh image, and the system halts, copies, and re-initializes the VMs securely."
            ],
            "image": "spotlight_vs_1.png"
        },
        {
            "title": "Logs",
            "text": ["The permanent system-of-record for every automated event in the platform."],
            "image": "screenshot_logs.png"
        },
        {
            "title": "Logs: Forensic Inspection",
            "text": [
                "Execution History Matrix: Grid of iterations tracking duration, targets, timestamps, and pass/fail outcomes.",
                "Deep Dive Inspection: Clicking a row pulls the raw console output trace for Engineering RCA.",
                "Export Utility: Download packages the raw logs safely for Jira/Emails."
            ],
            "image": "spotlight_logs_1.png"
        },
        {
            "title": "Terminal Integration",
            "text": ["Integrated Remote bash Web CLI bridges the gap without breaking workflow."],
            "image": "screenshot_terminal.png"
        },
        {
            "title": "Terminal: Direct Dispatch",
            "text": [
                "Direct Routing: Select target device to asynchronously broker a direct SSH tunnel.",
                "Command Dispatch: Type raw shell statements (ifconfig, show run). Pipeline returns near instantaneous results verifying switch state."
            ],
            "image": "spotlight_term_1.png"
        }
    ]

    for slide_data in slides_data:
        # Use Title + Content layout
        slide_layout = prs.slide_layouts[5] # Title only, so we can place text and image side by side manually
        slide = prs.slides.add_slide(slide_layout)
        
        shapes = slide.shapes
        shapes.title.text = slide_data['title']
        
        # Add textbox for bullets
        left = Inches(0.5)
        top = Inches(1.5)
        width = Inches(3.5)
        height = Inches(5.5)
        
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        
        for i, list_item in enumerate(slide_data['text']):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = list_item
            p.font.size = Pt(14)
            p.space_before = Pt(10)
            
        # Add image
        img_path = slide_data['image']
        if os.path.exists(img_path):
            try:
                # Place image on the right side
                left_img = Inches(4.2)
                top_img = Inches(1.5)
                # Keep aspect ratio by just providing width or height, maxing out at width 5.5
                slide.shapes.add_picture(img_path, left_img, top_img, width=Inches(5.5))
            except Exception as e:
                print(f"Error adding {img_path}: {e}")
        else:
            print(f"Warning: {img_path} not found.")

    try:
        prs.save("Eka_Automation_Demo.pptx")
        print("Presentation saved successfully to Eka_Automation_Demo.pptx")
    except Exception as e:
        print(f"Failed to save presentation: {e}")

if __name__ == "__main__":
    create_presentation()
