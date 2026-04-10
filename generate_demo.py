import os
import time
from playwright.sync_api import sync_playwright
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

def take_screenshots():
    screenshots = {}
    with sync_playwright() as p:
        browser = p.chromium.launch(headless=True)
        page = browser.new_page(viewport={"width": 1280, "height": 800})
        page.goto("http://localhost:8000/")
        time.sleep(2)  # Wait for load
        
        tabs = ['dashboard', 'devices', 'execute', 'vs', 'logs', 'terminal']
        for tab in tabs:
            page.evaluate(f"switchTab('{tab}')")
            time.sleep(1)
            filepath = f"screenshot_{tab}.png"
            page.screenshot(path=filepath)
            screenshots[tab] = filepath
            
        browser.close()
    return screenshots

def create_presentation(screenshots):
    prs = Presentation()
    
    # Title Slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Eka Automation Application"
    subtitle.text = "User Guide & Demo Presentation"
    
    # Slide 2: Agenda & Overview
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    shapes.title.text = "Agenda & Overview"
    body_shape = shapes.placeholders[1]
    tf = body_shape.text_frame
    tf.text = "What is Eka Automation?"
    p = tf.add_paragraph()
    p.text = "Application Navigation"
    p = tf.add_paragraph()
    p.text = "6 Core Tabs Walkthrough"
    
    # Slide 3: What is Eka Automation?
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    shapes.title.text = "What is Eka Automation?"
    body_shape = shapes.placeholders[1]
    tf = body_shape.text_frame
    tf.text = "Designed to significantly reduce repetitive workload through intelligent automation:"
    p = tf.add_paragraph()
    p.text = "Daily Image Loading - Automated batch processing of test images"
    p = tf.add_paragraph()
    p.text = "Batch Runs - Streamlined execution of multiple test configurations"
    p = tf.add_paragraph()
    p.text = "Advanced Regression Testing:"
    p.level = 0
    p2 = tf.add_paragraph()
    p2.text = "Parallel execution for faster test completion"
    p2.level = 1
    p3 = tf.add_paragraph()
    p3.text = "Optimized DUT (Device Under Test) usage efficiency"
    p3.level = 1
    p4 = tf.add_paragraph()
    p4.text = "Enhanced test coverage and reporting"
    p4.level = 1
    
    # Slide 4: Application Navigation
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    shapes.title.text = "Application Navigation - 6 Core Tabs"
    body_shape = shapes.placeholders[1]
    tf = body_shape.text_frame
    tf.text = "The application features a comprehensive interface organized into 6 main tabs:"
    for tab_name, desc in [
        ("Dashboard", "Overview and monitoring"),
        ("Devices", "Device management and configuration"),
        ("Execution", "Test execution and run management"),
        ("VS Manager", "Version and scenario management"),
        ("Logs", "System and test logging"),
        ("Terminal", "Direct command-line interface")
    ]:
        p = tf.add_paragraph()
        p.text = f"{tab_name} - {desc}"
        
    # Tab Slides with screenshots
    tab_details = {
        'dashboard': ("Dashboard Tab", "Overview of system status, total devices, online devices, and recent executions."),
        'devices': ("Devices Tab", "Manage testbed inventory. Add new VMs, DUTs, and Switches with IP and credentials."),
        'execute': ("Execution Tab", "Select VM, connect to Git repository, build Topology canvas, and execute test scripts."),
        'vs': ("VS Manager Tab", "Manage Virtual Machines. Apply new VS images, reset VMs, and check update progress."),
        'logs': ("Logs Tab", "View historical execution records and detailed real-time execution logs for debugging."),
        'terminal': ("Terminal Tab", "Access device command line remotely via SSH directly from the browser.")
    }
    
    for tab in ['dashboard', 'devices', 'execute', 'vs', 'logs', 'terminal']:
        slide_layout = prs.slide_layouts[5] # Title only
        slide = prs.slides.add_slide(slide_layout)
        shapes = slide.shapes
        title, desc = tab_details[tab]
        shapes.title.text = title
        
        # Add desc
        left = Inches(0.5)
        top = Inches(1.4)
        width = Inches(9)
        height = Inches(1)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = desc
        
        # Add image
        img_path = screenshots.get(tab)
        if img_path and os.path.exists(img_path):
            slide.shapes.add_picture(img_path, Inches(1), Inches(2.0), width=Inches(8))
            
    try:
        prs.save("Eka_Automation_Demo.pptx")
        print("Presentation saved to Eka_Automation_Demo.pptx")
    except Exception as e:
        print(f"Failed to save presentation: {e}")

if __name__ == "__main__":
    print("Taking screenshots...")
    screenshots = take_screenshots()
    print("Creating presentation...")
    create_presentation(screenshots)
