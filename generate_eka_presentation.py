import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

# Colors
BG_DARK = RGBColor(15, 17, 23)        # #0f1117 (Charcoal Dark)
BG_CARD = RGBColor(30, 41, 59)        # #1e293b (Slate Card Background)
BORDER_CARD = RGBColor(71, 85, 105)   # #475569 (Slate Border)
TEXT_WHITE = RGBColor(243, 244, 246)  # #f3f4f6 (Off-white)
TEXT_MUTED = RGBColor(156, 163, 175)  # #9ca3af (Gray Muted)
ACCENT_INDIGO = RGBColor(129, 140, 248) # #818cf8 (Indigo accent)
ACCENT_EMERALD = RGBColor(52, 211, 153) # #34d399 (Emerald Green)
ACCENT_RED = RGBColor(248, 113, 113)     # #f87171 (Red accent)

def set_slide_background(slide):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = BG_DARK

def add_title(slide, text):
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = 'Segoe UI'
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = ACCENT_INDIGO
    return title_box

def draw_card(slide, left, top, width, height, border_color=BORDER_CARD):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = BG_CARD
    card.line.color.rgb = border_color
    card.line.width = Pt(1.5)
    return card

def add_bullet_points(shape, title, points, title_color=ACCENT_INDIGO, body_color=TEXT_WHITE):
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.25)
    tf.margin_top = Inches(0.25)
    tf.margin_right = Inches(0.25)
    tf.margin_bottom = Inches(0.25)
    
    # Card Title
    p_title = tf.paragraphs[0]
    p_title.text = title
    p_title.font.name = 'Segoe UI'
    p_title.font.size = Pt(18)
    p_title.font.bold = True
    p_title.font.color.rgb = title_color
    p_title.space_after = Pt(12)
    
    # Bullets
    for pt in points:
        p = tf.add_paragraph()
        p.text = pt
        p.font.name = 'Calibri'
        p.font.size = Pt(13)
        p.font.color.rgb = body_color
        p.space_after = Pt(8)
        p.level = 0

def add_image_to_slide(slide, image_path, left, top, width):
    if os.path.exists(image_path):
        try:
            slide.shapes.add_picture(image_path, left, top, width=width)
            print(f"Added image {image_path} to slide.")
        except Exception as e:
            print(f"Error adding image {image_path}: {e}")
    else:
        # Draw placeholder
        rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Inches(4.5))
        rect.fill.solid()
        rect.fill.fore_color.rgb = BG_CARD
        rect.line.color.rgb = BORDER_CARD
        rect.line.width = Pt(1)
        tf = rect.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.text = f"\n\n\n[Image Placeholder]\n{image_path}\n(Run screenshots script to generate)"
        p.font.name = 'Segoe UI'
        p.font.size = Pt(12)
        p.font.color.rgb = TEXT_MUTED

def build_presentation():
    prs = Presentation()
    # Widescreen 16:9 layout
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # ----------------------------------------------------
    # Slide 1: Title Slide
    # ----------------------------------------------------
    slide = prs.slides.add_slide(prs.slide_layouts[6]) # blank layout
    set_slide_background(slide)
    
    # Left accent vertical bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(1.5), Inches(0.15), Inches(4.5))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT_INDIGO
    bar.line.fill.background()
    
    # Title Text Box
    title_box = slide.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(11.5), Inches(3.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "EKA AUTOMATION PLATFORM"
    p.font.name = 'Segoe UI'
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = TEXT_WHITE
    p.space_after = Pt(8)
    
    p2 = tf.add_paragraph()
    p2.text = "A Unified Control Plane for One-Click Network Simulation, Re-Imaging, and Test Validation"
    p2.font.name = 'Segoe UI'
    p2.font.size = Pt(18)
    p2.font.color.rgb = ACCENT_EMERALD
    p2.space_after = Pt(24)
    
    p3 = tf.add_paragraph()
    p3.text = "System Operations & Architecture Review"
    p3.font.name = 'Calibri'
    p3.font.size = Pt(14)
    p3.font.color.rgb = TEXT_MUTED

    # ----------------------------------------------------
    # Slide 2: Agenda
    # ----------------------------------------------------
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide)
    add_title(slide, "Presentation Agenda")
    
    card = draw_card(slide, Inches(0.5), Inches(1.5), Inches(6.5), Inches(5.2))
    points = [
        "1. What is Eka? — Sanskrit origin, vision, and core capabilities.",
        "2. Why Eka? — Operational pain points vs. Eka's integrated solutions.",
        "3. How Eka Works — Deep dive into direct switch re-imaging (Hardware Load) and virtual simulator patching (VS Manager).",
        "4. Interactive Topologies & Testing — Drag-and-drop cable mapping, Git integration, and concurrent regression execution.",
        "5. Modular Architecture — The 7-pillar subsystem framework.",
        "6. Interface Walkthrough — Tab-by-tab walkthrough of the live Eka web application with screenshots."
    ]
    add_bullet_points(card, "Core Topics Covered", points)
    
    # Agenda visual column (Right side decorative list summary)
    add_image_to_slide(slide, "screenshot_dashboard.png", Inches(7.3), Inches(1.8), Inches(5.5))

    # ----------------------------------------------------
    # Slide 3: What is Eka?
    # ----------------------------------------------------
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide)
    add_title(slide, "What is Eka?")
    
    card = draw_card(slide, Inches(0.5), Inches(1.5), Inches(6.0), Inches(5.2))
    points = [
        "Sanskrit Origin: Derived from Eka (एक), meaning 'Single' or 'One'. Symbolizes single-click, unified automation.",
        "Unified Automation Platform: Consolidates virtual simulation control, physical hardware re-imaging, cabling topology creation, and batch test execution.",
        "Centralized Dashboard: Provides real-time telemetry on devices, heartbeats, script syncs, and executions.",
        "One-Click Operations: Abstracts dozens of manual CLI syntax rules, complex directories, and credentials into simple UI actions."
    ]
    add_bullet_points(card, "Etymology & Core Purpose", points)
    
    add_image_to_slide(slide, "screenshot_dashboard.png", Inches(6.8), Inches(1.8), Inches(6.0))

    # ----------------------------------------------------
    # Slide 4: Why Eka? (Comparison Table)
    # ----------------------------------------------------
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide)
    add_title(slide, "Why Eka? Addressing Lab Inefficiencies")
    
    # Left Card - Challenges (Red Accent)
    card_left = draw_card(slide, Inches(0.5), Inches(1.5), Inches(6.0), Inches(5.2), border_color=ACCENT_RED)
    challenges = [
        "Manual Inefficiencies: Setting up VMs and flashing physical switches takes 30-40 minutes per engineer.",
        "Fragmentation: Operations require jumping between multiple Linux directories, TFTP servers, and Telnet tools.",
        "Syntax Barriers: Engineers must remember exact command syntax for different switches and ONIE bootloaders.",
        "Device Contention: Multiple users attempting to run tests on the same DUT, causing overlap conflicts."
    ]
    add_bullet_points(card_left, "Manual Operations (The Challenge)", challenges, title_color=ACCENT_RED)
    
    # Right Card - Eka Solutions (Emerald Accent)
    card_right = draw_card(slide, Inches(6.8), Inches(1.5), Inches(6.0), Inches(5.2), border_color=ACCENT_EMERALD)
    solutions = [
        "Instant Orchestration: Automated background imaging and VM patching triggered in seconds.",
        "Single Web Console: One interface handles inventory, topology drawing, terminal shells, and logs.",
        "Zero-CLI Dependency: Simple input forms collect IP/image parameters; Eka runs commands under the hood.",
        "Intelligent Locking: Integrates a central DUT reservation registry preventing overlapping execution."
    ]
    add_bullet_points(card_right, "Eka Platform (The Solution)", solutions, title_color=ACCENT_EMERALD)

    # ----------------------------------------------------
    # Slide 5: How Eka Works — Hardware Load
    # ----------------------------------------------------
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide)
    add_title(slide, "How Eka Works: Physical Hardware Load")
    
    card = draw_card(slide, Inches(0.5), Inches(1.5), Inches(6.0), Inches(5.2))
    points = [
        "Device Boot Detection: Connects via Telnet to identify if the switch is running a default image or stuck in an ONIE discovery loop.",
        "Discovery Interruption: Automatically issues discovery stop commands (e.g. 'onie-discovery-stop') to stabilize the shell.",
        "Automatic Config Scrubbing: Wipes existing startup configs to guarantee a clean build state.",
        "Integrated SCP Pipeline: Establishes a secure pipeline to copy the binary image from a developer server to the local switch filesystem.",
        "Automated Reboot: Issues reload sequence, tracks reboot steps, and prints live logs to the screen."
    ]
    add_bullet_points(card, "Physical Device Re-Imaging Pipeline", points)
    
    add_image_to_slide(slide, "screenshot_hardware_load.png", Inches(6.8), Inches(1.8), Inches(6.0))

    # ----------------------------------------------------
    # Slide 6: How Eka Works — VS Manager
    # ----------------------------------------------------
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide)
    add_title(slide, "How Eka Works: Virtual System (VS) Manager")
    
    card = draw_card(slide, Inches(0.5), Inches(1.5), Inches(6.0), Inches(5.2))
    points = [
        "Server Connection Layer: Connects securely to the hypervisor server hosting the virtual machines.",
        "Dynamic VM Discovery: Discovers all virtual simulator instances and lists their power status (Running, Stopped).",
        "Power Lifecycle Operations: Control virtual simulations using Start, Stop, and Reset buttons directly from the GUI.",
        "Robust SCP Implementation: Securely copies new virtual OS images (`.img`) from external build servers.",
        "Batch Virtual Provisioning: Apply a firmware image update to multiple VM nodes concurrently."
    ]
    add_bullet_points(card, "Virtual Simulation Orchestration", points)
    
    add_image_to_slide(slide, "screenshot_vs.png", Inches(6.8), Inches(1.8), Inches(6.0))

    # ----------------------------------------------------
    # Slide 7: How Eka Works — Batch Execution & Topology
    # ----------------------------------------------------
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide)
    add_title(slide, "How Eka Works: Execution & Topology Canvas")
    
    card = draw_card(slide, Inches(0.5), Inches(1.5), Inches(6.0), Inches(5.2))
    points = [
        "Interactive Topology Board: HTML5 canvas allows dragging and dropping of target nodes.",
        "Virtual Cable Routing: Visual port-to-port drawing link lines that auto-generate the master testbed YAML.",
        "Git Sync Integration: Sync and pull script repositories directly from the UI without manual command line git pull.",
        "SpyTest Batch Runner: Select scripts by category and queue them for execution based on device availability.",
        "Isolate Log Streaming: Streams logs live during execution and saves the final raw log file to the host VM."
    ]
    add_bullet_points(card, "Dynamic Cabling & Automation Dispatch", points)
    
    add_image_to_slide(slide, "screenshot_execute.png", Inches(6.8), Inches(1.8), Inches(6.0))

    # ----------------------------------------------------
    # Slide 8: Architecture Diagram
    # ----------------------------------------------------
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide)
    add_title(slide, "Eka Modular Platform: 7 Pillars")
    
    # Root Box (Eka Core Engine)
    root_shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.5), Inches(1.4), Inches(4.33), Inches(0.9))
    root_shape.fill.solid()
    root_shape.fill.fore_color.rgb = ACCENT_INDIGO
    root_shape.line.color.rgb = TEXT_WHITE
    root_shape.line.width = Pt(1.5)
    tf = root_shape.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = "EKA PLATFORM ENGINE\nFastAPI Backend | SQLite DB | HTML5 Frontend"
    p.font.name = 'Segoe UI'
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = BG_DARK
    
    # 7 Pillars arranged in two rows
    # Row 1 (4 pillars): Dashboard, Devices, Execution, VS Manager
    row1_tabs = [
        ("Dashboard", "Telemetry & Metrics", Inches(0.5)),
        ("Devices", "Inventory & heartbeats", Inches(3.7)),
        ("Execution", "Topology & SpyTest", Inches(6.9)),
        ("VS Manager", "VM Power & SCP Imaging", Inches(10.1))
    ]
    for name, desc, x in row1_tabs:
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(3.2), Inches(2.73), Inches(1.3))
        box.fill.solid()
        box.fill.fore_color.rgb = BG_CARD
        box.line.color.rgb = BORDER_CARD
        box.line.width = Pt(1)
        tf = box.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.text = name.upper()
        p.font.name = 'Segoe UI'
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = ACCENT_INDIGO
        
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        p2.text = desc
        p2.font.name = 'Calibri'
        p2.font.size = Pt(11)
        p2.font.color.rgb = TEXT_WHITE
        p2.space_before = Pt(4)
        
        # Add a simple line to root
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x + Inches(1.36), Inches(2.3), Inches(0.02), Inches(0.9))
        line.fill.solid()
        line.fill.fore_color.rgb = BORDER_CARD
        line.line.fill.background()

    # Row 2 (3 pillars): Logs, Terminal, Hardware Load
    row2_tabs = [
        ("Logs", "Forensic Run Logs", Inches(2.1)),
        ("Terminal", "SSH Browser Shells", Inches(5.3)),
        ("Hardware Load", "Physical Switch Imaging", Inches(8.5))
    ]
    for name, desc, x in row2_tabs:
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(5.2), Inches(2.73), Inches(1.3))
        box.fill.solid()
        box.fill.fore_color.rgb = BG_CARD
        box.line.color.rgb = BORDER_CARD
        box.line.width = Pt(1)
        tf = box.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.text = name.upper()
        p.font.name = 'Segoe UI'
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = ACCENT_INDIGO
        
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        p2.text = desc
        p2.font.name = 'Calibri'
        p2.font.size = Pt(11)
        p2.font.color.rgb = TEXT_WHITE
        p2.space_before = Pt(4)
        
        # Add a simple line to root region
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x + Inches(1.36), Inches(4.5), Inches(0.02), Inches(0.7))
        line.fill.solid()
        line.fill.fore_color.rgb = BORDER_CARD
        line.line.fill.background()
        
    # Draw horizontal crossbar
    crossbar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.86), Inches(2.3), Inches(9.6), Inches(0.02))
    crossbar.fill.solid()
    crossbar.fill.fore_color.rgb = BORDER_CARD
    crossbar.line.fill.background()

    # ----------------------------------------------------
    # Slides 9-15: Tab-by-tab walkthrough
    # ----------------------------------------------------
    tabs_data = [
        {
            "id": "1",
            "title": "Dashboard Telemetry",
            "desc": "The central operations console showing overall infrastructure status.",
            "points": [
                "Infrastructure Pulse: Real-time telemetry widgets showing total registered lab devices.",
                "Reachability Status: Live statistics on online vs. offline devices checked by background pings.",
                "System Counters: Lifelong counters of test executions and indexed Git repository scripts.",
                "Clean Visual Design: Uses custom HSL variables to style cards in a sleek dark theme."
            ],
            "img": "screenshot_dashboard.png",
            "spotlight": "spotlight_dashboard_1.png"
        },
        {
            "id": "2",
            "title": "Devices Registry",
            "desc": "Manage the registered switch inventory and virtual simulators.",
            "points": [
                "Simplified Onboarding: Standardized form to ingest name, IP, port, type, and login details.",
                "Secure Key Management: Stores passwords safely and encrypts credentials at rest.",
                "On-Demand Validation: Wi-Fi symbol action button triggers immediate background ping checks.",
                "Device Clean-out: Clear trash-can action triggers safe purge of decommissioned device records."
            ],
            "img": "screenshot_devices.png",
            "spotlight": "spotlight_devices_2.png"
        },
        {
            "id": "3",
            "title": "Execution Framework",
            "desc": "Map network layouts, sync code repositories, and dispatch scripts.",
            "points": [
                "Draggable Cable Designer: Connect nodes visually to construct network links.",
                "Git Synchronization: Clones repositories and indexes categories/scripts dynamically in the backend.",
                "Testbed Automation: Select the target YAML testbed and automatically compile a master topology config.",
                "Parallel SpyTest Execution: Dispatches regression scripts concurrently on VM targets."
            ],
            "img": "screenshot_execute.png",
            "spotlight": "spotlight_execute_2.png"
        },
        {
            "id": "4",
            "title": "VS Manager Console",
            "desc": "Orchestrates version imaging and power lifecycle of virtual simulations.",
            "points": [
                "Auto-Discovery of Instances: Queries the host server to list all active Virtual Machines.",
                "Direct Power Controls: Toggle VMs immediately with Start, Stop, and Reset action buttons.",
                "SCP Image Synchronization: Transfers sonic-vs image binaries directly from dev servers.",
                "Batch Execution: Select multiple virtual instances to re-image them concurrently."
            ],
            "img": "screenshot_vs.png",
            "spotlight": "spotlight_vs_1.png"
        },
        {
            "id": "5",
            "title": "Forensic Run Logs",
            "desc": "Detailed record of past automated events and diagnostic consoles.",
            "points": [
                "Audit Logging Matrix: Central logs grid containing timestamps, durations, and outcomes.",
                "RCA Deep-Dive Inspector: Select an execution row to pull the raw console output logs.",
                "Log Retention: Allows clearing history and cleaning outdated execution logs.",
                "Export Capability: Packages raw execution logs for rapid engineering feedback."
            ],
            "img": "screenshot_logs.png",
            "spotlight": "spotlight_logs_1.png"
        },
        {
            "id": "6",
            "title": "SSH Web Terminal",
            "desc": "Browser-based CLI tool matching the utility of standalone Putty.",
            "points": [
                "Direct Shell Tunneling: Select any registered device IP to broker an instant remote command-line session.",
                "Terminal Emulation: Renders terminal scripts (e.g. show version, config interfaces) immediately.",
                "Multiplexed Sessions: Uses tmux or screen under the hood to preserve active terminal sessions.",
                "Integrated Workflow: Eliminates the need to jump between external terminal clients."
            ],
            "img": "screenshot_terminal.png",
            "spotlight": "spotlight_term_1.png"
        },
        {
            "id": "7",
            "title": "Hardware Load Pipeline",
            "desc": "Full-cycle automated switch flash and boot loader management.",
            "points": [
                "ONIE Intercept: Auto-detects physical hardware state and runs discovery stop sequences.",
                "Config Scrubbing: Wipes existing startup configs to guarantee a clean switch status.",
                "Integrated SCP Flash: Securely transfers binary firmware builds and flashes the physical ASIC.",
                "Active Progress Audit: Live visual progress bars and terminal logger tracking the installation lifecycle.",
                "Record Cleanup: Add and delete execution job logs to maintain a clean database."
            ],
            "img": "screenshot_hardware_load.png",
            "spotlight": "spotlight_hardware_load_1.png"
        }
    ]
    
    for tdata in tabs_data:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        set_slide_background(slide)
        
        # Heading
        add_title(slide, f"{tdata['id']}. {tdata['title']}")
        
        # Left Details Card
        card = draw_card(slide, Inches(0.5), Inches(1.5), Inches(5.5), Inches(5.2))
        add_bullet_points(card, tdata['desc'], tdata['points'])
        
        # Right Images (Main screenshot + highlighted spotlight overlapping/aligned)
        # Main screenshot scaled down
        add_image_to_slide(slide, tdata['img'], Inches(6.3), Inches(1.5), Inches(6.5))
        
        # Spotlight insert at bottom right if available
        if tdata['spotlight'] and os.path.exists(tdata['spotlight']):
            try:
                # Overlap at bottom right
                slide.shapes.add_picture(tdata['spotlight'], Inches(8.5), Inches(3.9), width=Inches(4.3))
                print(f"Added spotlight overlay {tdata['spotlight']}.")
            except Exception as e:
                print(f"Error adding spotlight {tdata['spotlight']}: {e}")

    # ----------------------------------------------------
    # Slide 16: Summary & ROI
    # ----------------------------------------------------
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide)
    add_title(slide, "Eka platform: Value & Efficiency Summary")
    
    # 3 horizontal columns
    cols = [
        ("TIME & SPEED SAVINGS", 
         [
             "Cuts setup time from 40 mins to 1 click.",
             "Automated file transfers and imaging run concurrently in the background.",
             "Enables zero-touch regression overnight."
         ], Inches(0.5), ACCENT_INDIGO),
         
        ("INTEGRATION & COVERAGE", 
         [
             "Draw and cable topologies in HTML5 canvas.",
             "Sync Git and trigger SpyTest regression batches directly.",
             "Bridges virtual testing (VS) and physical lab switches (Hardware)."
         ], Inches(4.7), ACCENT_EMERALD),
         
        ("LAB UTILIZATION & ROI", 
         [
             "DUT reservation locks prevent collision.",
             "Saves engineers from memorizing commands.",
             "Optimizes hardware scheduling and maximizes test coverages."
         ], Inches(8.9), ACCENT_EMERALD)
    ]
    
    for title, points, x, accent_color in cols:
        card = draw_card(slide, x, Inches(1.8), Inches(3.9), Inches(4.8), border_color=accent_color)
        add_bullet_points(card, title, points, title_color=accent_color)

    # Save
    filename = "Eka_Automation_Presentation.pptx"
    try:
        prs.save(filename)
        print(f"\n=======================================================")
        print(f"[Success] PowerPoint saved to {filename}")
        print(f"=======================================================")
    except Exception as e:
        print(f"Failed to save PPTX: {e}")

if __name__ == "__main__":
    build_presentation()
