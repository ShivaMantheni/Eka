import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from PIL import Image

# ----------------------------------------------------
# NATURAL & ARTISTIC GALLERY DESIGN SYSTEM - COLOR PALETTE
# ----------------------------------------------------
# A warm, modern design aesthetic inspired by high-end design portfolios and clean tech UI.
BG_ALABASTER = RGBColor(247, 245, 240)      # #f7f5f0 (Organic paper feel)
BG_WHITE = RGBColor(255, 255, 255)          # #ffffff (Crisp cotton/linen card)
BORDER_CLAY = RGBColor(226, 220, 208)        # #e2dcd0 (Subtle clay line)
TEXT_CHARCOAL = RGBColor(26, 36, 28)        # #1a241c (Spruce charcoal black)
TEXT_MUTED = RGBColor(94, 104, 95)          # #5e685f (Olive slate gray)
ACCENT_GREEN = RGBColor(36, 77, 50)          # #244d32 (Forest pine green)
ACCENT_TERRACOTTA = RGBColor(176, 89, 62)    # #b0593e (Warm copper terracotta)

# Simplified and polished slide content (beginner-friendly but grammatically perfect)
agenda_points = [
    "What is Eka? (Origin, Vision, and Core Concept)",
    "Why Eka? (Simplifying Laboratory Complexity)",
    "How Eka Works (Provisioning, Virtual Systems, and Executions)",
    "Modular Architecture (The 7 Core Operational Tabs)"
]

what_is_eka_points = [
    "Eka is an automation platform designed to consolidate virtual simulation loading, physical hardware re-imaging, and test script execution into a single-click workflow.",
    "Named after the Sanskrit word 'Eka' (एक) meaning 'Single' or 'One' — representing a unified core interface."
]

why_eka_points = [
    "Traditional virtual system and hardware setups take several minutes and involve a complex sequence of CLI commands, credentials, and directories.",
    "Eka abstracts these complexities, removing the need to remember manual commands or file paths.",
    "Streamlines parallel test execution and automated regressions using SpyTest.",
    "Saves engineering hours and maximizes the efficiency of shared physical lab hardware."
]

how_hardware_points = [
    "Users can update physical switches with new build images simply by entering the device IP, credentials, and image path.",
    "Eka automatically detects whether the switch has a running operating system or is in a raw bootloader state (ONIE loop).",
    "It performs a clean image installation, removes legacy configuration files, and restarts the device safely."
]

how_vs_points = [
    "Users enter hypervisor details to discover and display all running virtual systems and VMs in a central view.",
    "Start or stop virtual systems instantly with interactive dashboard toggles.",
    "Integrates a built-in SCP workflow to fetch and copy build images from any reachable network server.",
    "Allows engineers to update single or multiple virtual systems simultaneously in a batch."
]

how_execution_points = [
    "Browse virtual machine folders and select script files directly from the browser window.",
    "Drag and drop devices onto the interactive canvas to draw topology lines and connect ports.",
    "Automatically compiles and runs SpyTest execution queues, recording console logs directly in the VM.",
    "Supports Git pull integrations to sync test suites, and runs multiple tests concurrently."
]

# Tab Walkthrough configurations using the newly created mockup images
tabs_data = [
    {
        "id": "1",
        "title": "Dashboard",
        "desc": "High-level overview of testbeds and active runs.",
        "points": [
            "Monitor all onboarded devices and check their live connection status.",
            "Track active script execution progress and see overall run counters.",
            "Display metrics showing the count of online, offline, and running systems."
        ],
        "img": "mockup_dashboard.png"
    },
    {
        "id": "2",
        "title": "Devices",
        "desc": "Inventory management and live reachability.",
        "points": [
            "Onboard new physical switches and virtual devices into the network list.",
            "Edit registration configurations such as IP addresses, login credentials, and types.",
            "Perform instant background ping checks by clicking the action Wi-Fi icon."
        ],
        "img": "mockup_devices.png"
    },
    {
        "id": "3",
        "title": "Execution",
        "desc": "Drag-and-drop testbed canvas and script dispatcher.",
        "points": [
            "Execute automated regression scripts directly from the interface.",
            "Draw virtual cables between device ports on an interactive drag-and-drop canvas.",
            "Follow active scripts in real time using the live terminal execution stream."
        ],
        "img": "mockup_execute.png"
    },
    {
        "id": "4",
        "title": "VS Manager",
        "desc": "Virtual system management and SCP imaging.",
        "points": [
            "View and manage running virtual machines from a single hypervisor query.",
            "Start or stop simulation nodes with simple dashboard power controls.",
            "Batch update VM images using SCP transfer from reachable development hosts."
        ],
        "img": "mockup_vs.png"
    },
    {
        "id": "5",
        "title": "Logs",
        "desc": "Audit trail database and diagnostic log viewer.",
        "points": [
            "Access a detailed historical timeline of all completed test executions.",
            "Select individual execution records to view detailed console log traces.",
            "Retrieve and review raw diagnostic files to troubleshoot test failures."
        ],
        "img": "mockup_logs.png"
    },
    {
        "id": "6",
        "title": "Terminal",
        "desc": "Integrated PuTTY-style SSH console utility.",
        "points": [
            "Select any registered device to instantly open an SSH session in the browser.",
            "Maintains session configurations similar to standard terminal clients like PuTTY.",
            "Supports background terminal persistence using tmux or screen sessions."
        ],
        "img": "mockup_terminal.png"
    },
    {
        "id": "7",
        "title": "Hardware Load",
        "desc": "Physical switch installer and re-imaging loop control.",
        "points": [
            "Provide device IP, credentials, and the firmware path to kick off re-imaging.",
            "Detects and clean-updates switches that are empty or have existing builds.",
            "Monitor background installation logs directly inside the workflow panel."
        ],
        "img": "mockup_hardware_load.png"
    }
]

# ----------------------------------------------------
# 1. GENERATE PPTX PRESENTATION (Segoe UI/Calibri Modern Clean Style)
# ----------------------------------------------------
def set_slide_background(slide, prs):
    # Draw a full screen rectangle to guarantee Alabaster background color
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG_ALABASTER
    bg.line.color.rgb = BG_ALABASTER
    bg.line.width = Pt(0)

def add_title(slide, text):
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = 'Segoe UI'
    p.font.size = Pt(30)
    p.font.bold = True
    p.font.color.rgb = ACCENT_GREEN
    return title_box

def draw_card(slide, left, top, width, height, border_color=BORDER_CLAY):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = BG_WHITE
    card.line.color.rgb = border_color
    card.line.width = Pt(1.2)
    return card

def add_bullet_points(shape, title, points, title_color=ACCENT_TERRACOTTA, body_color=TEXT_CHARCOAL):
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.3)
    tf.margin_top = Inches(0.3)
    tf.margin_right = Inches(0.3)
    tf.margin_bottom = Inches(0.3)
    
    p_title = tf.paragraphs[0]
    p_title.text = title
    p_title.font.name = 'Segoe UI'
    p_title.font.size = Pt(20)
    p_title.font.bold = True
    p_title.font.color.rgb = title_color
    p_title.space_after = Pt(14)
    
    for pt in points:
        p = tf.add_paragraph()
        p.text = "—  " + pt
        p.font.name = 'Calibri'
        p.font.size = Pt(14)
        p.font.color.rgb = body_color
        p.space_after = Pt(10)
        p.level = 0

def add_image_to_slide(slide, image_path, left, top, width):
    if os.path.exists(image_path):
        try:
            # Dynamic aspect ratio loading to size the gallery mat board exactly
            img = Image.open(image_path)
            w_px, h_px = img.size
            aspect = h_px / w_px
            height = width * aspect
            
            # Mat dimensions
            mat_width = width + Inches(0.2)
            mat_height = height + Inches(0.2)
            
            # Draw the clean white mat board with clay border
            mat = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left - Inches(0.1), top - Inches(0.1), mat_width, mat_height)
            mat.fill.solid()
            mat.fill.fore_color.rgb = BG_WHITE
            mat.line.color.rgb = BORDER_CLAY
            mat.line.width = Pt(1)
            
            # Place the picture exactly in the center of the mat
            slide.shapes.add_picture(image_path, left, top, width=width)
        except Exception as e:
            print(f"Error adding image {image_path}: {e}")
    else:
        # Drawing a clean placeholder card
        rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Inches(4.2))
        rect.fill.solid()
        rect.fill.fore_color.rgb = BG_WHITE
        rect.line.color.rgb = BORDER_CLAY
        rect.line.width = Pt(1)
        tf = rect.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.text = f"\n\n\n[Gallery UI Mockup]\n{image_path}"
        p.font.name = 'Segoe UI'
        p.font.size = Pt(12)
        p.font.color.rgb = TEXT_MUTED

def compile_pptx():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Slide 1: Title Slide (Warm, Clean, and Elegant)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, prs)
    
    # Left organic accent border
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.0), Inches(2.0), Inches(0.08), Inches(3.5))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT_GREEN
    bar.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(1.4), Inches(2.1), Inches(10.5), Inches(3.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "EKA AUTOMATION PLATFORM"
    p.font.name = 'Segoe UI'
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = ACCENT_GREEN
    p.space_after = Pt(8)
    
    p2 = tf.add_paragraph()
    p2.text = "A Unified System for One-Click Lab Simulation and Re-Imaging"
    p2.font.name = 'Segoe UI'
    p2.font.size = Pt(16)
    p2.font.bold = True
    p2.font.color.rgb = ACCENT_TERRACOTTA
    
    p3 = tf.add_paragraph()
    p3.text = "System Operations & Architecture Review"
    p3.font.name = 'Calibri'
    p3.font.size = Pt(11)
    p3.font.color.rgb = TEXT_MUTED
    p3.space_before = Pt(36)
    
    # Slide 2: Agenda
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, prs)
    add_title(slide, "Agenda")
    card = draw_card(slide, Inches(0.5), Inches(1.5), Inches(6.2), Inches(5.2))
    add_bullet_points(card, "Operational Objectives", agenda_points)
    add_image_to_slide(slide, "mockup_dashboard.png", Inches(7.3), Inches(1.8), Inches(5.5))
    
    # Slide 3: What is Eka?
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, prs)
    add_title(slide, "What is Eka?")
    card = draw_card(slide, Inches(0.5), Inches(1.5), Inches(6.0), Inches(5.2))
    add_bullet_points(card, "Eka (एक) • Unified Control Plane", what_is_eka_points)
    add_image_to_slide(slide, "mockup_dashboard.png", Inches(7.0), Inches(1.8), Inches(5.8))
    
    # Slide 4: Why Eka?
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, prs)
    add_title(slide, "Why Eka?")
    card = draw_card(slide, Inches(0.5), Inches(1.5), Inches(12.33), Inches(5.2))
    add_bullet_points(card, "Simplifying Laboratory Complexity", why_eka_points, title_color=ACCENT_GREEN)
    
    # Slide 5: How Eka works - Hardware Load
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, prs)
    add_title(slide, "How Eka works - Hardware Load")
    card = draw_card(slide, Inches(0.5), Inches(1.5), Inches(6.0), Inches(5.2))
    add_bullet_points(card, "Physical Device Re-Imaging", how_hardware_points)
    add_image_to_slide(slide, "mockup_hardware_load.png", Inches(7.0), Inches(1.8), Inches(5.8))
    
    # Slide 6: How Eka works - VS Manager
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, prs)
    add_title(slide, "How Eka works - VS Manager")
    card = draw_card(slide, Inches(0.5), Inches(1.5), Inches(6.0), Inches(5.2))
    add_bullet_points(card, "Virtual Simulation Manager", how_vs_points)
    add_image_to_slide(slide, "mockup_vs.png", Inches(7.0), Inches(1.8), Inches(5.8))
    
    # Slide 7: How Eka works - Script Execution
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, prs)
    add_title(slide, "How Eka works - Script Execution")
    card = draw_card(slide, Inches(0.5), Inches(1.5), Inches(6.0), Inches(5.2))
    add_bullet_points(card, "Integrated Topology & Execution", how_execution_points)
    add_image_to_slide(slide, "mockup_execute.png", Inches(7.0), Inches(1.8), Inches(5.8))

    # Slide 8: Architecture Diagram (Clean, Architectural Style)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, prs)
    add_title(slide, "Modular System Architecture")
    
    # Root Core Plaque
    root_shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.5), Inches(1.4), Inches(4.33), Inches(0.9))
    root_shape.fill.solid()
    root_shape.fill.fore_color.rgb = ACCENT_GREEN
    root_shape.line.color.rgb = BG_ALABASTER
    root_shape.line.width = Pt(1.5)
    tf = root_shape.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = "EKA AUTOMATION PLATFORM\n7 Core Subsystem Modules"
    p.font.name = 'Segoe UI'
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = BG_ALABASTER
    
    # Connecting Lines in background color
    crossbar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.86), Inches(2.3), Inches(9.6), Inches(0.02))
    crossbar.fill.solid()
    crossbar.fill.fore_color.rgb = BORDER_CLAY
    crossbar.line.fill.background()

    row1_tabs = [
        ("Dashboard", "See onboarded, running, online status.", Inches(0.5)),
        ("Devices", "Onboard device and edit details.", Inches(3.7)),
        ("Execution", "Run scripts and create testbed.", Inches(6.9)),
        ("VS Manager", "VM power and image updates.", Inches(10.1))
    ]
    for name, desc, x in row1_tabs:
        # Pillar Card
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(3.2), Inches(2.73), Inches(1.3))
        box.fill.solid()
        box.fill.fore_color.rgb = BG_WHITE
        box.line.color.rgb = BORDER_CLAY
        box.line.width = Pt(1)
        tf = box.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.text = name.upper()
        p.font.name = 'Segoe UI'
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = ACCENT_TERRACOTTA
        
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        p2.text = desc
        p2.font.name = 'Calibri'
        p2.font.size = Pt(10)
        p2.font.color.rgb = TEXT_CHARCOAL
        p2.space_before = Pt(4)
        
        # Link line to central crossbar
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x + Inches(1.36), Inches(2.3), Inches(0.02), Inches(0.9))
        line.fill.solid()
        line.fill.fore_color.rgb = BORDER_CLAY
        line.line.fill.background()

    row2_tabs = [
        ("Logs", "Completed logs history list.", Inches(2.1)),
        ("Terminal", "Configure device like PuTTY.", Inches(5.3)),
        ("Hardware Load", "Image physical switch hardware.", Inches(8.5))
    ]
    for name, desc, x in row2_tabs:
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(5.2), Inches(2.73), Inches(1.3))
        box.fill.solid()
        box.fill.fore_color.rgb = BG_WHITE
        box.line.color.rgb = BORDER_CLAY
        box.line.width = Pt(1)
        tf = box.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.text = name.upper()
        p.font.name = 'Segoe UI'
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = ACCENT_TERRACOTTA
        
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        p2.text = desc
        p2.font.name = 'Calibri'
        p2.font.size = Pt(10)
        p2.font.color.rgb = TEXT_CHARCOAL
        p2.space_before = Pt(4)
        
        # Link line to central crossbar
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x + Inches(1.36), Inches(4.5), Inches(0.02), Inches(0.7))
        line.fill.solid()
        line.fill.fore_color.rgb = BORDER_CLAY
        line.line.fill.background()

    # Slides 9-15: Operational Tabs walkthrough
    for tdata in tabs_data:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        set_slide_background(slide, prs)
        
        add_title(slide, f"{tdata['id']}. {tdata['title']} Tab")
        card = draw_card(slide, Inches(0.5), Inches(1.5), Inches(5.5), Inches(5.2))
        add_bullet_points(card, tdata['desc'], tdata['points'])
        
        # Centered mockup image frame (no overlapping spotlights, yielding clean asymmetrical layouts)
        add_image_to_slide(slide, tdata['img'], Inches(6.4), Inches(1.8), Inches(6.4))

    # Slide 16: Summary
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, prs)
    add_title(slide, "Summary")
    card = draw_card(slide, Inches(0.5), Inches(1.8), Inches(12.33), Inches(4.8))
    summary_bullets = [
        "Saves Valuable Engineering Time: Removes memory blocks, complex CLI command pathways, and folder directory loops.",
        "Reduces Fragmented Labor: Automation pools re-image switches and start virtual server machines seamlessly.",
        "Maximizes Lab Resource Allocation: Automated locking prevents resource collision and coordinates execution queues."
    ]
    add_bullet_points(card, "Eka Platform Value Propositions", summary_bullets, title_color=ACCENT_GREEN)

    filename = "Eka_Automation_Presentation.pptx"
    prs.save(filename)
    print(f"[OK] Saved PPTX to {filename}")


# ----------------------------------------------------
# 2. GENERATE REVEAL.JS HTML PRESENTATION (Outfit/Plus Jakarta Sans Modern Style)
# ----------------------------------------------------
def compile_html():
    html_content = """<!doctype html>
<html lang="en">
<head>
    <meta charset="utf-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0, maximum-scale=1.0, user-scalable=no">

    <title>Eka Automation Platform Presentation</title>

    <link rel="stylesheet" href="https://cdnjs.cloudflare.com/ajax/libs/reveal.js/5.1.0/reset.min.css">
    <link rel="stylesheet" href="https://cdnjs.cloudflare.com/ajax/libs/reveal.js/5.1.0/reveal.min.css">
    <link rel="stylesheet" href="https://cdnjs.cloudflare.com/ajax/libs/reveal.js/5.1.0/theme/serif.min.css">
    <link href="https://fonts.googleapis.com/icon?family=Material+Icons+Round" rel="stylesheet">
    <link href="https://fonts.googleapis.com/css2?family=Outfit:wght@400;500;600;700&family=Plus+Jakarta+Sans:wght@400;500;600;700&display=swap" rel="stylesheet">

    <style>
        body {
            background-color: #f7f5f0 !important; /* Warm Alabaster */
            font-family: 'Plus Jakarta Sans', sans-serif;
            color: #1a241c; /* Spruce Charcoal */
        }
        .reveal {
            font-family: 'Plus Jakarta Sans', sans-serif;
            color: #1a241c;
        }
        .reveal h1, .reveal h2, .reveal h3, .reveal h4 {
            font-family: 'Outfit', sans-serif;
            text-transform: none;
            color: #244d32; /* Spruce Green */
            font-weight: 700;
            text-shadow: none;
            letter-spacing: -0.02em;
        }
        .reveal h1 {
            font-size: 2.6em;
            line-height: 1.1;
            margin-bottom: 8px;
        }
        .reveal h2 {
            color: #244d32;
            font-size: 1.65em;
            margin-bottom: 24px;
            border-bottom: 1px double #e2dcd0;
            padding-bottom: 12px;
            text-align: left;
        }
        .slide-card {
            background: #ffffff;
            border: 1px solid #e2dcd0;
            border-radius: 12px;
            padding: 24px 30px;
            box-shadow: 0 10px 30px rgba(45, 38, 25, 0.04);
            text-align: left;
            margin-top: 10px;
        }
        .card-title {
            color: #b0593e; /* Terracotta copper */
            font-family: 'Outfit', sans-serif;
            font-size: 1.05em;
            font-weight: 700;
            margin-bottom: 18px;
            display: flex;
            align-items: center;
            gap: 10px;
            border-bottom: 1px dashed rgba(226, 220, 208, 0.6);
            padding-bottom: 8px;
            letter-spacing: 0.5px;
            text-transform: uppercase;
        }
        .reveal ul {
            font-size: 0.62em;
            line-height: 1.6;
            list-style-type: none;
            margin: 0;
            padding: 0;
        }
        .reveal li {
            margin-bottom: 12px;
            position: relative;
            padding-left: 22px;
            color: #4a544b;
        }
        .reveal li::before {
            content: "—";
            color: #b0593e;
            font-weight: bold;
            position: absolute;
            left: 0;
            top: -1px;
        }
        .two-column {
            display: flex;
            flex-direction: row;
            align-items: center;
            gap: 35px;
            width: 100%;
        }
        .col-left {
            flex: 1;
            display: flex;
            flex-direction: column;
            justify-content: flex-start;
        }
        .col-right {
            flex: 1.1;
            position: relative;
            display: flex;
            align-items: center;
            justify-content: center;
        }
        /* Matted Gallery Frame style for images */
        .gallery-frame {
            background: #ffffff;
            border: 1px solid #dcd6c8;
            padding: 14px;
            box-shadow: 0 12px 32px rgba(45, 38, 25, 0.07);
            border-radius: 4px;
            width: 100%;
            display: block;
        }
        .image-frame {
            width: 100%;
            border-radius: 2px;
            border: 1px solid #e2dcd0;
            background: #fafaf9;
            display: block;
        }
        .sanskrit-box {
            font-family: 'Outfit', sans-serif;
            font-size: 0.95em;
            color: #244d32;
            font-weight: 700;
            text-transform: uppercase;
            letter-spacing: 0.5px;
            background: rgba(36, 77, 50, 0.05);
            padding: 6px 16px;
            border-radius: 6px;
            display: inline-block;
            margin-bottom: 16px;
            border: 1px solid rgba(36, 77, 50, 0.1);
        }
        .arch-container {
            display: flex;
            flex-direction: column;
            align-items: center;
            gap: 20px;
            width: 100%;
            margin-top: 10px;
        }
        .arch-engine {
            background: #ffffff;
            color: #244d32;
            padding: 14px 28px;
            border-radius: 6px;
            font-weight: 700;
            font-family: 'Outfit', sans-serif;
            font-size: 0.75em;
            box-shadow: 0 8px 24px rgba(45, 38, 25, 0.06);
            border: 1px solid #244d32;
            text-align: center;
            letter-spacing: 0.5px;
        }
        .arch-grid-row1 {
            display: grid;
            grid-template-columns: repeat(4, 1fr);
            gap: 15px;
            width: 100%;
        }
        .arch-grid-row2 {
            display: flex;
            justify-content: center;
            gap: 15px;
            width: 100%;
        }
        .arch-pillar {
            background: #ffffff;
            border: 1px solid #e2dcd0;
            border-radius: 8px;
            padding: 12px;
            text-align: center;
            box-shadow: 0 4px 15px rgba(45, 38, 25, 0.02);
            transition: all 0.4s ease;
        }
        .arch-pillar:hover {
            border-color: #b0593e;
            transform: translateY(-2px);
            box-shadow: 0 10px 25px rgba(176, 89, 62, 0.08);
        }
        .arch-pillar-title {
            color: #b0593e;
            font-weight: 700;
            font-size: 0.58em;
            margin-bottom: 6px;
            font-family: 'Outfit', sans-serif;
            letter-spacing: 0.5px;
        }
        .arch-pillar-desc {
            color: #5e685f;
            font-size: 0.42em;
            line-height: 1.4;
        }
        .arch-line {
            width: 1px;
            height: 24px;
            background: #e2dcd0;
        }
        .nav-helper {
            font-size: 0.38em !important;
            color: #5e685f;
            font-style: italic;
            text-align: center;
            margin-top: 20px;
        }
        .reveal .controls {
            color: #b0593e !important;
        }
        .reveal .progress {
            color: #b0593e !important;
            background: rgba(226, 220, 208, 0.4) !important;
        }
        .reveal .progress span {
            background: #b0593e !important;
        }
        .download-btn {
            display: inline-flex;
            align-items: center;
            gap: 8px;
            background: #244d32;
            color: #ffffff !important;
            padding: 10px 20px;
            border-radius: 6px;
            font-size: 0.50em !important;
            font-family: 'Outfit', sans-serif;
            font-weight: 600;
            text-decoration: none;
            box-shadow: 0 4px 15px rgba(36, 77, 50, 0.15);
            transition: all 0.2s ease;
            margin-top: 20px;
            border: 1px solid rgba(255, 255, 255, 0.1);
        }
        .download-btn:hover {
            background: #1b3a26;
            transform: translateY(-1px);
            box-shadow: 0 6px 20px rgba(36, 77, 50, 0.25);
        }
    </style>
</head>
<body>
    <div class="reveal">
        <div class="slides">

            <!-- Slide 1: Title -->
            <section style="text-align: left;">
                <div style="display: flex; gap: 40px; align-items: stretch; margin-left: 20px;">
                    <div style="width: 4px; background: #244d32; border-radius: 2px;"></div>
                    <div>
                        <h1 style="margin: 0; padding: 0; font-size: 2.5em; font-weight: 700;">EKA AUTOMATION PLATFORM</h1>
                        <p style="color: #b0593e; font-size: 0.8em; font-family: 'Outfit', sans-serif; font-weight: 600; margin: 15px 0 25px 0;">
                            A Unified System for One-Click Lab Simulation and Re-Imaging
                        </p>
                        <p style="color: #5e685f; font-size: 0.45em; letter-spacing: 1px; text-transform: uppercase; margin-bottom: 25px;">System Operations & Architecture Review</p>
                        
                        <a href="Eka_Automation_Presentation.pptx" download class="download-btn">
                            <span class="material-icons-round" style="font-size: 18px;">download</span> 
                            Download PowerPoint (.pptx)
                        </a>
                    </div>
                </div>
            </section>

            <!-- Slide 2: Agenda -->
            <section>
                <h2>Agenda</h2>
                <div class="two-column">
                    <div class="col-left">
                        <div class="slide-card" style="height: 100%;">
                            <div class="card-title"><span class="material-icons-round" style="color: #b0593e;">menu_book</span> Operational Objectives</div>
                            <ul>
                                <li><strong>What is Eka?</strong> (Origin, Vision, and Core Concept)</li>
                                <li><strong>Why Eka?</strong> (Simplifying Laboratory Complexity)</li>
                                <li><strong>How Eka Works</strong> (Provisioning, Virtual Systems, and Executions)</li>
                                <li><strong>Modular Architecture</strong> (The 7 Core Operational Tabs)</li>
                            </ul>
                        </div>
                    </div>
                    <div class="col-right">
                        <div class="gallery-frame">
                            <img src="mockup_dashboard.png" class="image-frame" alt="Eka Dashboard">
                        </div>
                    </div>
                </div>
            </section>

            <!-- Slide 3: What is Eka? -->
            <section>
                <h2>What is Eka?</h2>
                <div class="two-column">
                    <div class="col-left">
                        <div class="slide-card" style="height: 100%;">
                            <div class="sanskrit-box">एक (Eka) &bull; Unified Control Plane</div>
                            <div class="card-title"><span class="material-icons-round" style="color: #b0593e;">info</span> Core Vision</div>
                            <ul>
                                <li>Eka is an automation platform designed to consolidate virtual simulation loading, physical hardware re-imaging, and test script execution into a single-click workflow.</li>
                                <li>Named after the Sanskrit word 'Eka' (एक) meaning 'Single' or 'One' — representing a unified core interface.</li>
                            </ul>
                        </div>
                    </div>
                    <div class="col-right">
                        <div class="gallery-frame">
                            <img src="mockup_dashboard.png" class="image-frame" alt="Eka Dashboard">
                        </div>
                    </div>
                </div>
            </section>

            <!-- Slide 4: Why Eka? -->
            <section>
                <h2>Why Eka?</h2>
                <div class="two-column">
                    <div class="col-left" style="flex: 1.4;">
                        <div class="slide-card" style="height: 100%;">
                            <div class="card-title"><span class="material-icons-round" style="color: #b0593e;">help_outline</span> Addressing Laboratory Complexity</div>
                            <ul>
                                <li>Traditional virtual system and hardware setups take several minutes and involve a complex sequence of CLI commands, credentials, and directories.</li>
                                <li>Eka abstracts these complexities, removing the need to remember manual commands or file paths.</li>
                                <li>Streamlines parallel test execution and automated regressions using SpyTest.</li>
                                <li>Saves engineering hours and maximizes the efficiency of shared physical lab hardware.</li>
                            </ul>
                        </div>
                    </div>
                    <div class="col-right" style="flex: 0.9;">
                        <div class="gallery-frame">
                            <img src="mockup_dashboard.png" class="image-frame" alt="Dashboard">
                        </div>
                    </div>
                </div>
            </section>

            <!-- Slide 5: How Eka Works - Hardware Load -->
            <section>
                <h2>How Eka Works - Hardware Load</h2>
                <div class="two-column">
                    <div class="col-left">
                        <div class="slide-card" style="height: 100%;">
                            <div class="card-title"><span class="material-icons-round" style="color: #b0593e;">settings_input_hdmi</span> Physical Device Re-Imaging</div>
                            <ul>
                                <li>Users can update physical switches with new build images simply by entering the device IP, credentials, and image path.</li>
                                <li>Eka automatically detects whether the switch has a running operating system or is in a raw bootloader state (ONIE loop).</li>
                                <li>It performs a clean image installation, removes legacy configuration files, and restarts the device safely.</li>
                            </ul>
                        </div>
                    </div>
                    <div class="col-right">
                        <div class="gallery-frame">
                            <img src="mockup_hardware_load.png" class="image-frame" alt="Hardware Load">
                        </div>
                    </div>
                </div>
            </section>

            <!-- Slide 6: How Eka Works - VS Manager -->
            <section>
                <h2>How Eka Works - VS Manager</h2>
                <div class="two-column">
                    <div class="col-left">
                        <div class="slide-card" style="height: 100%;">
                            <div class="card-title"><span class="material-icons-round" style="color: #b0593e;">grid_view</span> Virtual Simulation Manager</div>
                            <ul>
                                <li>Users enter hypervisor details to discover and display all running virtual systems and VMs in a central view.</li>
                                <li>Start or stop virtual systems instantly with interactive dashboard toggles.</li>
                                <li>Integrates a built-in SCP workflow to fetch and copy build images from any reachable network server.</li>
                                <li>Allows engineers to update single or multiple virtual systems simultaneously in a batch.</li>
                            </ul>
                        </div>
                    </div>
                    <div class="col-right">
                        <div class="gallery-frame">
                            <img src="mockup_vs.png" class="image-frame" alt="VS Manager">
                        </div>
                    </div>
                </div>
            </section>

            <!-- Slide 7: How Eka Works - Script Execution -->
            <section>
                <h2>How Eka Works - Script Execution</h2>
                <div class="two-column">
                    <div class="col-left">
                        <div class="slide-card" style="height: 100%;">
                            <div class="card-title"><span class="material-icons-round" style="color: #b0593e;">play_circle</span> Integrated Topology & Execution</div>
                            <ul>
                                <li>Browse virtual machine folders and select script files directly from the browser window.</li>
                                <li>Drag and drop devices onto the interactive canvas to draw topology lines and connect ports.</li>
                                <li>Automatically compiles and runs SpyTest execution queues, recording console logs directly in the VM.</li>
                                <li>Supports Git pull integrations to sync test suites, and runs multiple tests concurrently.</li>
                            </ul>
                        </div>
                    </div>
                    <div class="col-right">
                        <div class="gallery-frame">
                            <img src="mockup_execute.png" class="image-frame" alt="Script Execution">
                        </div>
                    </div>
                </div>
            </section>

            <!-- Slide 8: Architecture -->
            <section>
                <h2>Modular System Architecture</h2>
                <div class="arch-container">
                    <div class="arch-engine">
                        EKA AUTOMATION PLATFORM
                        <div style="font-size: 0.5em; font-weight: normal; margin-top: 4px; color: #5e685f;">
                            Consolidated Subsystem Control Plane
                        </div>
                    </div>
                    
                    <div class="arch-line"></div>
                    
                    <div class="arch-grid-row1">
                        <div class="arch-pillar">
                            <div class="arch-pillar-title">DASHBOARD</div>
                            <div class="arch-pillar-desc">Monitor onboarded devices, active executions, and metrics.</div>
                        </div>
                        <div class="arch-pillar">
                            <div class="arch-pillar-title">DEVICES</div>
                            <div class="arch-pillar-desc">Inventory management, credentials, and reachability.</div>
                        </div>
                        <div class="arch-pillar">
                            <div class="arch-pillar-title">EXECUTION</div>
                            <div class="arch-pillar-desc">Script selection, topology canvas, and parallel execution.</div>
                        </div>
                        <div class="arch-pillar">
                            <div class="arch-pillar-title">VS MANAGER</div>
                            <div class="arch-pillar-desc">Virtual setup state control and remote SCP updates.</div>
                        </div>
                    </div>
                    
                    <div class="arch-grid-row2">
                        <div class="arch-pillar" style="width: 25%;">
                            <div class="arch-pillar-title">LOGS</div>
                            <div class="arch-pillar-desc">Execution database logs and diagnostic file browser.</div>
                        </div>
                        <div class="arch-pillar" style="width: 25%;">
                            <div class="arch-pillar-title">TERMINAL</div>
                            <div class="arch-pillar-desc">Direct PuTTY-style remote SSH shells in browser.</div>
                        </div>
                        <div class="arch-pillar" style="width: 25%;">
                            <div class="arch-pillar-title">HARDWARE LOAD</div>
                            <div class="arch-pillar-desc">Physical switch provisioning and bootloader loop control.</div>
                        </div>
                    </div>
                </div>
                <div class="nav-helper">Use keyboard arrows to navigate. Press ESC for layout overview.</div>
            </section>

            <!-- Slides 9-15: Walkthrough of 7 Tabs -->
            <!-- Dashboard Tab -->
            <section>
                <h2>1. Dashboard Tab</h2>
                <div class="two-column">
                    <div class="col-left">
                        <div class="slide-card" style="height: 100%;">
                            <div class="card-title"><span class="material-icons-round" style="color: #b0593e;">dashboard</span> System Overview</div>
                            <ul>
                                <li>Monitor all onboarded devices and check their live connection status.</li>
                                <li>Track active script execution progress and see overall run counters.</li>
                                <li>Display metrics showing the count of online, offline, and running systems.</li>
                            </ul>
                        </div>
                    </div>
                    <div class="col-right">
                        <div class="gallery-frame">
                            <img src="mockup_dashboard.png" class="image-frame" alt="Dashboard">
                        </div>
                    </div>
                </div>
            </section>

            <!-- Devices Tab -->
            <section>
                <h2>2. Devices Tab</h2>
                <div class="two-column">
                    <div class="col-left">
                        <div class="slide-card" style="height: 100%;">
                            <div class="card-title"><span class="material-icons-round" style="color: #b0593e;">devices_other</span> Inventory Inventory</div>
                            <ul>
                                <li>Onboard new physical switches and virtual devices into the network list.</li>
                                <li>Edit registration configurations such as IP addresses, login credentials, and types.</li>
                                <li>Perform instant background ping checks by clicking the action Wi-Fi icon.</li>
                            </ul>
                        </div>
                    </div>
                    <div class="col-right">
                        <div class="gallery-frame">
                            <img src="mockup_devices.png" class="image-frame" alt="Devices">
                        </div>
                    </div>
                </div>
            </section>

            <!-- Execution Tab -->
            <section>
                <h2>3. Execution Tab</h2>
                <div class="two-column">
                    <div class="col-left">
                        <div class="slide-card" style="height: 100%;">
                            <div class="card-title"><span class="material-icons-round" style="color: #b0593e;">play_circle</span> Interactive Canvas</div>
                            <ul>
                                <li>Execute automated regression scripts directly from the interface.</li>
                                <li>Draw virtual cables between device ports on an interactive drag-and-drop canvas.</li>
                                <li>Follow active scripts in real time using the live terminal execution stream.</li>
                            </ul>
                        </div>
                    </div>
                    <div class="col-right">
                        <div class="gallery-frame">
                            <img src="mockup_execute.png" class="image-frame" alt="Execution">
                        </div>
                    </div>
                </div>
            </section>

            <!-- VS Manager Tab -->
            <section>
                <h2>4. VS Manager Tab</h2>
                <div class="two-column">
                    <div class="col-left">
                        <div class="slide-card" style="height: 100%;">
                            <div class="card-title"><span class="material-icons-round" style="color: #b0593e;">memory</span> VM Power & Control</div>
                            <ul>
                                <li>View and manage running virtual machines from a single hypervisor query.</li>
                                <li>Start or stop simulation nodes with simple dashboard power controls.</li>
                                <li>Batch update VM images using SCP transfer from reachable development hosts.</li>
                            </ul>
                        </div>
                    </div>
                    <div class="col-right">
                        <div class="gallery-frame">
                            <img src="mockup_vs.png" class="image-frame" alt="VS Manager">
                        </div>
                    </div>
                </div>
            </section>

            <!-- Logs Tab -->
            <section>
                <h2>5. Logs Tab</h2>
                <div class="two-column">
                    <div class="col-left">
                        <div class="slide-card" style="height: 100%;">
                            <div class="card-title"><span class="material-icons-round" style="color: #b0593e;">article</span> Audit Trail & Debugging</div>
                            <ul>
                                <li>Access a detailed historical timeline of all completed test executions.</li>
                                <li>Select individual execution records to view detailed console log traces.</li>
                                <li>Retrieve and review raw diagnostic files to troubleshoot test failures.</li>
                            </ul>
                        </div>
                    </div>
                    <div class="col-right">
                        <div class="gallery-frame">
                            <img src="mockup_logs.png" class="image-frame" alt="Logs">
                        </div>
                    </div>
                </div>
            </section>

            <!-- Terminal Tab -->
            <section>
                <h2>6. Terminal Tab</h2>
                <div class="two-column">
                    <div class="col-left">
                        <div class="slide-card" style="height: 100%;">
                            <div class="card-title"><span class="material-icons-round" style="color: #b0593e;">terminal</span> Browser SSH Console</div>
                            <ul>
                                <li>Select any registered device to instantly open an SSH session in the browser.</li>
                                <li>Maintains session configurations similar to standard terminal clients like PuTTY.</li>
                                <li>Supports background terminal persistence using tmux or screen sessions.</li>
                            </ul>
                        </div>
                    </div>
                    <div class="col-right">
                        <div class="gallery-frame">
                            <img src="mockup_terminal.png" class="image-frame" alt="Terminal">
                        </div>
                    </div>
                </div>
            </section>

            <!-- Hardware Load Tab -->
            <section>
                <h2>7. Hardware Load Tab</h2>
                <div class="two-column">
                    <div class="col-left">
                        <div class="slide-card" style="height: 100%;">
                            <div class="card-title"><span class="material-icons-round" style="color: #b0593e;">settings_input_hdmi</span> Re-Imaging Loop Control</div>
                            <ul>
                                <li>Provide device IP, credentials, and the firmware path to kick off re-imaging.</li>
                                <li>Detects and clean-updates switches that are empty or have existing builds.</li>
                                <li>Monitor background installation logs directly inside the workflow panel.</li>
                            </ul>
                        </div>
                    </div>
                    <div class="col-right">
                        <div class="gallery-frame">
                            <img src="mockup_hardware_load.png" class="image-frame" alt="Hardware Load">
                        </div>
                    </div>
                </div>
            </section>

            <!-- Slide 16: Summary -->
            <section style="text-align: left;">
                <h2>Summary</h2>
                <div class="slide-card">
                    <div class="card-title" style="color: #244d32;"><span class="material-icons-round">check_circle</span> Value Proposition Summary</div>
                    <ul>
                        <li><strong>Saves Valuable Engineering Time:</strong> Removes memory blocks, complex CLI command pathways, and folder directory loops.</li>
                        <li><strong>Reduces Fragmented Labor:</strong> Automation pools re-image switches and start virtual server machines seamlessly.</li>
                        <li><strong>Maximizes Lab Resource Allocation:</strong> Automated locking prevents resource collision and coordinates execution queues.</li>
                    </ul>
                </div>
            </section>

            <!-- End Slide -->
            <section style="text-align: center;">
                <h1 style="font-size: 3em;">Thank You.</h1>
                <p style="color: #b0593e; font-size: 1.1em; font-family: 'Outfit', sans-serif; font-weight: 600;">Ready for Automated Lab Validation</p>
                <div class="nav-helper">Press ESC to see slide overview. Use arrow keys to navigate.</div>
            </section>

        </div>
    </div>

    <script src="https://cdnjs.cloudflare.com/ajax/libs/reveal.js/5.1.0/reveal.min.js"></script>
    <script>
        Reveal.initialize({
            hash: true,
            slideNumber: 'c/t',
            progress: true,
            center: true,
            width: 1280,
            height: 720,
            transition: 'fade',
            controls: true,
            keyboard: true
        });
    </script>
</body>
</html>
"""
    # Make sure static directory exists
    os.makedirs("static", exist_ok=True)

    with open("Eka_Automation_Web_Presentation.html", "w", encoding="utf-8") as f:
        f.write(html_content)
    print("[OK] Saved HTML to Eka_Automation_Web_Presentation.html")

    # Write to static directory for direct HTTP serving
    with open("static/Eka_Automation_Web_Presentation.html", "w", encoding="utf-8") as f:
        f.write(html_content)
    print("[OK] Saved HTML to static/Eka_Automation_Web_Presentation.html")

    # Copy all mockup images to static directory
    import shutil
    images_to_copy = [
        "mockup_dashboard.png", "mockup_devices.png", "mockup_execute.png", 
        "mockup_vs.png", "mockup_logs.png", "mockup_terminal.png", "mockup_hardware_load.png"
    ]
    for img in images_to_copy:
        if os.path.exists(img):
            shutil.copy(img, os.path.join("static", img))
            print(f"[OK] Copied {img} to static/")
            
    # Copy PPTX to static folder for direct browser download
    if os.path.exists("Eka_Automation_Presentation.pptx"):
        shutil.copy("Eka_Automation_Presentation.pptx", os.path.join("static", "Eka_Automation_Presentation.pptx"))
        print("[OK] Copied Eka_Automation_Presentation.pptx to static/")


if __name__ == "__main__":
    print("Compiling PPTX (Modern Clean Style)...")
    compile_pptx()
    print("Compiling HTML (Modern Clean Style)...")
    compile_html()
    print("Done compiling presentations successfully!")
