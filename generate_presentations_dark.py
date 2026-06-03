import os
import shutil
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from PIL import Image

# ----------------------------------------------------
# FUTURISTIC CYBER-NETWORKING DARK PALETTE
# ----------------------------------------------------
BG_DARK = RGBColor(10, 16, 26)          # Deep space dark navy (#0a101a)
BG_CARD = RGBColor(18, 28, 45)          # Glassmorphism card fill (#121c2d)
BORDER_CYAN = RGBColor(0, 168, 204)     # Neon cyan border (#00a8cc)
BORDER_DARK = RGBColor(26, 47, 76)      # Slate blue border (#1a2f4c)
TEXT_WHITE = RGBColor(255, 255, 255)    # High-contrast white (#ffffff)
TEXT_MUTED = RGBColor(136, 153, 184)    # Muted ice slate (#8899b8)
ACCENT_CYAN = RGBColor(0, 240, 255)     # Glowing electric cyan (#00f0ff)
ACCENT_BLUE = RGBColor(0, 122, 255)     # Bright blue accent (#007aff)

# Slide agenda points & details
agenda_points = [
    "What is Eka? (Origin, Vision, and Core Control Plane)",
    "Problem Statement (Laboratory Complexity & Manual Roadblocks)",
    "Solution Overview (Unified Simulation & Imaging Engine)",
    "Application Architecture (7 Core Subsystem Modules)",
    "Key Features (Canvas Topology, ONIE loops, SSH persistent)",
    "Workflow & Automation Pipeline (Step-by-step Run Process)",
    "Benefits (Maximizing Engineering Hours & Shared Hardware)",
    "Demo & Tab Walkthrough (Live App Screenshots Show)",
    "Roadmap & Conclusion (Future SSO & Distributed Test Grids)"
]

# Tab Walkthrough configurations referencing high-res screenshots
tabs_data = [
    {
        "id": "1",
        "title": "Dashboard Tab",
        "desc": "High-Level Control Panel overview of Eka instances.",
        "points": [
            "Monitor live connection states of all registered devices in real time.",
            "Display visual counters for total devices, online nodes, and script queues.",
            "Review run outcomes via recent execution telemetry charts."
        ],
        "img": "screenshot_dashboard.png"
    },
    {
        "id": "2",
        "title": "Devices Tab",
        "desc": "Inventory management and live ping reachability checker.",
        "points": [
            "Onboard physical switches and virtual machines with IP credentials.",
            "Modify network parameter mappings (Ports, SSH credentials, paths).",
            "Perform background ping heartbeats to determine device active states."
        ],
        "img": "screenshot_devices.png"
    },
    {
        "id": "3",
        "title": "Execution Tab",
        "desc": "Drag-and-Drop Topology builder and automation pipeline launcher.",
        "points": [
            "Draw cabling connections between virtual device interfaces on a canvas.",
            "Browse hierarchical test script folders and select test groups.",
            "Trigger automated SpyTest executions with a single click."
        ],
        "img": "screenshot_execute.png"
    },
    {
        "id": "4",
        "title": "VS Manager Tab",
        "desc": "Virtual System power toggling and VM image updating.",
        "points": [
            "Query and display VM states from remote hypervisors.",
            "Power on/off individual or groups of virtual switches instantly.",
            "Fetch and copy VM build images from network storage via SCP."
        ],
        "img": "screenshot_vs.png"
    },
    {
        "id": "5",
        "title": "Logs Tab",
        "desc": "Historical execution database and terminal trace logger.",
        "points": [
            "Access database audit logs for all completed and failed runs.",
            "Expand execution rows to inspect live terminal stdio streams.",
            "Filter records by date, script name, or test engineer."
        ],
        "img": "screenshot_logs.png"
    },
    {
        "id": "6",
        "title": "Terminal Tab",
        "desc": "Integrated multi-session SSH console tool (PuTTY-style).",
        "points": [
            "Open instant SSH terminal sessions to any registered device inside the UI.",
            "Execute standard CLI commands without leaving the web browser application.",
            "Preserve console states using background tmux session persistence."
        ],
        "img": "screenshot_terminal.png"
    },
    {
        "id": "7",
        "title": "Hardware Load Tab",
        "desc": "Physical switch firmware re-imaging loop controller.",
        "points": [
            "Update switch builds automatically via a single-button installer.",
            "Auto-detect switches in ONIE boot loops and install firmware packages.",
            "View real-time console boot sequence logs during re-imaging."
        ],
        "img": "screenshot_hardware_load.png"
    }
]

def set_slide_background(slide, prs):
    # Draw a solid background rectangle in dark space navy
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG_DARK
    bg.line.color.rgb = BG_DARK
    
    # Add a thin neon cyan glowing line at the top border (glowing accent)
    accent_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, Inches(0.06))
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = ACCENT_CYAN
    accent_bar.line.fill.background()

def add_title(slide, text):
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = 'Segoe UI'
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = ACCENT_CYAN
    return title_box

def draw_card(slide, left, top, width, height, border_color=BORDER_DARK):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = BG_CARD
    card.line.color.rgb = border_color
    card.line.width = Pt(1.5)
    return card

def add_bullet_points(shape, title, points, title_color=ACCENT_CYAN, body_color=TEXT_WHITE):
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.3)
    tf.margin_top = Inches(0.3)
    tf.margin_right = Inches(0.3)
    tf.margin_bottom = Inches(0.3)
    
    p_title = tf.paragraphs[0]
    p_title.text = title
    p_title.font.name = 'Segoe UI'
    p_title.font.size = Pt(22)
    p_title.font.bold = True
    p_title.font.color.rgb = title_color
    p_title.space_after = Pt(16)
    
    for pt in points:
        p = tf.add_paragraph()
        p.text = "✦  " + pt
        p.font.name = 'Calibri'
        p.font.size = Pt(15)
        p.font.color.rgb = body_color
        p.space_after = Pt(12)
        p.level = 0

def add_image_to_slide(slide, image_path, left, top, width):
    if os.path.exists(image_path):
        try:
            # Aspect ratio check
            img = Image.open(image_path)
            w_px, h_px = img.size
            aspect = h_px / w_px
            height = width * aspect
            
            # Draw card board back frame (glowing mat)
            mat_w = width + Inches(0.12)
            mat_h = height + Inches(0.12)
            mat = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left - Inches(0.06), top - Inches(0.06), mat_w, mat_h)
            mat.fill.solid()
            mat.fill.fore_color.rgb = BG_CARD
            mat.line.color.rgb = BORDER_CYAN
            mat.line.width = Pt(1.5)
            
            # Place screenshot
            slide.shapes.add_picture(image_path, left, top, width=width)
        except Exception as e:
            print(f"Error adding image {image_path}: {e}")
    else:
        # Placeholder card
        rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Inches(4.2))
        rect.fill.solid()
        rect.fill.fore_color.rgb = BG_CARD
        rect.line.color.rgb = BORDER_CYAN
        rect.line.width = Pt(1.5)
        tf = rect.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.text = f"\n\n\n[Eka Screenshot Space]\n{image_path}"
        p.font.name = 'Segoe UI'
        p.font.size = Pt(14)
        p.font.color.rgb = TEXT_MUTED

def compile_pptx():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # ----------------------------------------------------
    # SLIDE 1: Title Slide (Futuristic Dark Theme)
    # ----------------------------------------------------
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, prs)
    
    # Draw topology decorative background lines
    l1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.5), Inches(2.2), Inches(0.06), Inches(3.2))
    l1.fill.solid()
    l1.fill.fore_color.rgb = ACCENT_CYAN
    l1.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(1.8), Inches(2.1), Inches(10.5), Inches(3.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "EKA AUTOMATION PLATFORM"
    p.font.name = 'Segoe UI'
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = TEXT_WHITE
    p.space_after = Pt(6)
    
    p2 = tf.add_paragraph()
    p2.text = "Enterprise-Grade One-Click Network Testbeds & Simulation Control Plane"
    p2.font.name = 'Segoe UI'
    p2.font.size = Pt(18)
    p2.font.color.rgb = ACCENT_CYAN
    p2.space_after = Pt(28)
    
    p3 = tf.add_paragraph()
    p3.text = "System Operations & Architecture Review  •  For Management and Engineering"
    p3.font.name = 'Calibri'
    p3.font.size = Pt(12)
    p3.font.color.rgb = TEXT_MUTED
    
    # ----------------------------------------------------
    # SLIDE 2: Agenda
    # ----------------------------------------------------
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, prs)
    add_title(slide, "Presentation Agenda")
    card = draw_card(slide, Inches(0.5), Inches(1.5), Inches(6.0), Inches(5.2))
    add_bullet_points(card, "Operational Core Elements", agenda_points)
    add_image_to_slide(slide, "screenshot_dashboard.png", Inches(7.0), Inches(1.8), Inches(5.8))
    
    # ----------------------------------------------------
    # SLIDE 3: Problem Statement
    # ----------------------------------------------------
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, prs)
    add_title(slide, "Problem Statement: Testing & Dev Friction")
    card = draw_card(slide, Inches(0.5), Inches(1.5), Inches(12.33), Inches(5.2))
    problems = [
        "Time-Consuming Lab Setup: Traditional virtual systems and physical hardware staging take hours of complex command-line entry.",
        "Error-Prone Manual Procedures: Accessing hypervisors, setting credentials, and copying firmware manually often lead to configuration drift.",
        "Fragmented Tooling: Developers use separate terminal tools (PuTTY), file copy tools (SCP), and custom testing modules (SpyTest).",
        "Resource Collision: Multiple engineers attempting to connect or re-image the same switches simultaneously without central coordination."
    ]
    add_bullet_points(card, "Laboratory Complexity Roadblocks", problems, title_color=ACCENT_CYAN)
    
    # ----------------------------------------------------
    # SLIDE 4: Solution Overview
    # ----------------------------------------------------
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, prs)
    add_title(slide, "Solution Overview: Unified Automation")
    card = draw_card(slide, Inches(0.5), Inches(1.5), Inches(12.33), Inches(5.2))
    solutions = [
        "Unified Control Plane: Eka aggregates virtual machines, physical switches, and execution frameworks into one web interface.",
        "One-Click Image Deployment: Automated routines copy build packages via SCP and re-image firmware safely.",
        "Asynchronous Execution Pipelines: Runs testbeds concurrently and streams test result stdout logs directly to active panels.",
        "Resource Isolation & Heartbeats: Background monitors block overlapping DUT reservation sessions automatically."
    ]
    add_bullet_points(card, "The Eka Platform Answer", solutions, title_color=ACCENT_CYAN)
    
    # ----------------------------------------------------
    # SLIDE 5: Application Architecture
    # ----------------------------------------------------
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, prs)
    add_title(slide, "Modular System Architecture")
    
    # Architecture Central block
    root_shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.5), Inches(1.4), Inches(4.33), Inches(0.9))
    root_shape.fill.solid()
    root_shape.fill.fore_color.rgb = ACCENT_CYAN
    root_shape.line.color.rgb = TEXT_WHITE
    root_shape.line.width = Pt(1.5)
    tf = root_shape.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = "EKA CONTROL ENGINE\nUnified Orchestration Interface"
    p.font.name = 'Segoe UI'
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = BG_DARK
    
    # Connector line background
    crossbar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.86), Inches(2.3), Inches(9.6), Inches(0.02))
    crossbar.fill.solid()
    crossbar.fill.fore_color.rgb = BORDER_CYAN
    crossbar.line.fill.background()
    
    # Submodules Row 1
    row1_tabs = [
        ("Dashboard", "Dynamic stats visualizer.", Inches(0.5)),
        ("Devices", "Centralized switch registry.", Inches(3.7)),
        ("Execution", "Topology drag-and-drop canvas.", Inches(6.9)),
        ("VS Manager", "VM hypervisor coordinator.", Inches(10.1))
    ]
    for name, desc, x in row1_tabs:
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(3.2), Inches(2.73), Inches(1.3))
        box.fill.solid()
        box.fill.fore_color.rgb = BG_CARD
        box.line.color.rgb = BORDER_DARK
        box.line.width = Pt(1.5)
        tf = box.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.text = name.upper()
        p.font.name = 'Segoe UI'
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = ACCENT_CYAN
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        p2.text = desc
        p2.font.name = 'Calibri'
        p2.font.size = Pt(10)
        p2.font.color.rgb = TEXT_WHITE
        
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x + Inches(1.36), Inches(2.3), Inches(0.02), Inches(0.9))
        line.fill.solid()
        line.fill.fore_color.rgb = BORDER_CYAN
        line.line.fill.background()
        
    # Submodules Row 2
    row2_tabs = [
        ("Logs & Audits", "Consolidated trace file DB.", Inches(2.1)),
        ("Terminal SSH", "Integrated PuTTY web client.", Inches(5.3)),
        ("Hardware Load", "ONIE re-imaging pipeline.", Inches(8.5))
    ]
    for name, desc, x in row2_tabs:
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(5.2), Inches(2.73), Inches(1.3))
        box.fill.solid()
        box.fill.fore_color.rgb = BG_CARD
        box.line.color.rgb = BORDER_DARK
        box.line.width = Pt(1.5)
        tf = box.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.text = name.upper()
        p.font.name = 'Segoe UI'
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = ACCENT_CYAN
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        p2.text = desc
        p2.font.name = 'Calibri'
        p2.font.size = Pt(10)
        p2.font.color.rgb = TEXT_WHITE
        
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x + Inches(1.36), Inches(4.5), Inches(0.02), Inches(0.7))
        line.fill.solid()
        line.fill.fore_color.rgb = BORDER_CYAN
        line.line.fill.background()
        
    # ----------------------------------------------------
    # SLIDE 6: Key Features
    # ----------------------------------------------------
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, prs)
    add_title(slide, "Eka Core Feature Highlights")
    card = draw_card(slide, Inches(0.5), Inches(1.5), Inches(12.33), Inches(5.2))
    features = [
        "Dynamic Topology Builder: Check checklist boxes and drag ports to connect virtual Ethernet interfaces instantly.",
        "ONIE Re-imaging Loop Control: Monitors switches in zero-OS state and flashes custom firmware asynchronously.",
        "Persistent Web SSH: Open active switch CLI terminals via browser websockets; continues running scripts inside tmux.",
        "Hypervisor Integration: Instantly check VMs list, power state, IPs, and update target VM builds.",
        "Database Audit Trails: Stores execution duration, exit states, and standard logs securely for team reviews."
    ]
    add_bullet_points(card, "Powering Modern Test Automation", features, title_color=ACCENT_CYAN)
    
    # ----------------------------------------------------
    # SLIDE 7: Workflow/Automation Pipeline
    # ----------------------------------------------------
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, prs)
    add_title(slide, "Automation Pipeline Pipeline")
    
    # Add pipeline nodes as chevron shapes or boxes linked together
    steps = [
        ("Step 1: Onboard", "Register target VM or switch IP credentials in Devices panel."),
        ("Step 2: Provision", "Execute re-imaging loops or hypervisor build updates via SCP."),
        ("Step 3: Connect", "Draw virtual interface cable topologies on active canvas."),
        ("Step 4: Execute", "Launch test files and trace stdout outputs dynamically."),
        ("Step 5: Review", "Audit terminal output history in consolidated logs database.")
    ]
    for i, (stitle, sdesc) in enumerate(steps):
        y = 1.6 + (i * 1.05)
        # Glow dot
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.8), Inches(y + 0.1), Inches(0.3), Inches(0.3))
        dot.fill.solid()
        dot.fill.fore_color.rgb = ACCENT_CYAN
        dot.line.color.rgb = TEXT_WHITE
        
        # Link line to next dot (except last)
        if i < 4:
            line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.94), Inches(y + 0.4), Inches(0.02), Inches(0.8))
            line.fill.solid()
            line.fill.fore_color.rgb = BORDER_DARK
            line.line.fill.background()
            
        tb = slide.shapes.add_textbox(Inches(1.3), Inches(y), Inches(11.2), Inches(0.9))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = stitle
        p.font.name = 'Segoe UI'
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = ACCENT_CYAN
        p2 = tf.add_paragraph()
        p2.text = sdesc
        p2.font.name = 'Calibri'
        p2.font.size = Pt(13)
        p2.font.color.rgb = TEXT_WHITE
        
    # ----------------------------------------------------
    # SLIDE 8: Benefits for Developers & Testers
    # ----------------------------------------------------
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, prs)
    add_title(slide, "Platform Benefits: Dev & QA Teams")
    card = draw_card(slide, Inches(0.5), Inches(1.5), Inches(12.33), Inches(5.2))
    benefits = [
        "Time Saved: Automated ONIE loops and hypervisor flash workflows save up to 45 minutes per switch setup cycle.",
        "Zero Access Overhead: Access devices from any location via browser websockets; no VPN credentials or client apps (PuTTY) needed.",
        "Enhanced Reliability: Eliminates configuration drift by loading pristine build packages and clearing system storage before runs.",
        "Optimized Hardware Usage: Automated locking prevents device collision; schedules regression tasks dynamically."
    ]
    add_bullet_points(card, "Engineering Advantages", benefits, title_color=ACCENT_CYAN)
    
    # ----------------------------------------------------
    # SLIDE 9: Scalability and Performance
    # ----------------------------------------------------
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, prs)
    add_title(slide, "Scalability & Performance Metrics")
    card = draw_card(slide, Inches(0.5), Inches(1.5), Inches(12.33), Inches(5.2))
    metrics = [
        "Parallel Installation Loops: Multi-threaded installer flash operations run on up to 8 switches concurrently.",
        "Lightweight Event Bus: FastAPI async background workers handle connection polling with less than 2% hypervisor overhead.",
        "Optimized SSH Pooling: Paramiko connection pools maintain background sessions, dropping connection latency by 85%.",
        "Horizontal Architecture: The server logic can be containerized using Docker and scaled to support dozens of concurrent users."
    ]
    add_bullet_points(card, "Built for Enterprise scale", metrics, title_color=ACCENT_CYAN)
    
    # ----------------------------------------------------
    # SLIDES 10 - 16: Demo & Walkthrough (Actual Screenshots)
    # ----------------------------------------------------
    for tdata in tabs_data:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        set_slide_background(slide, prs)
        add_title(slide, f"Demo: {tdata['title']}")
        
        card = draw_card(slide, Inches(0.5), Inches(1.5), Inches(5.5), Inches(5.2))
        add_bullet_points(card, tdata['desc'], tdata['points'])
        
        # Display 3x crisp image on right
        add_image_to_slide(slide, tdata['img'], Inches(6.3), Inches(1.6), Inches(6.5))
        
    # ----------------------------------------------------
    # SLIDE 17: Conclusion & Future Roadmap
    # ----------------------------------------------------
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, prs)
    add_title(slide, "Roadmap & Future Extensions")
    card = draw_card(slide, Inches(0.5), Inches(1.5), Inches(12.33), Inches(5.2))
    roadmap = [
        "OIDC Single-Sign-On Integration: Bypassing local user tables to delegate auth queries to corporate enterprise portals (OnePalc).",
        "Distributed Test Grids: Virtual cable topo canvas to coordinate test runs on remote multi-hypervisor nodes.",
        "AI Log Analysis: Integrating diagnostic models to parse execution files and automatically isolate network failures.",
        "Hardware Health Telemetry: Dashboard charts displaying live temperature, fan load, and memory usage for test switches."
    ]
    add_bullet_points(card, "Platform Roadmap", roadmap, title_color=ACCENT_CYAN)
    
    filename = "Eka_Automation_Presentation.pptx"
    prs.save(filename)
    print(f"[OK] Saved PPTX to {filename}")

# ----------------------------------------------------
# GENERATE REVEAL.JS HTML PRESENTATION
# ----------------------------------------------------
def compile_html():
    html_content = """<!doctype html>
<html lang="en">
<head>
    <meta charset="utf-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0, maximum-scale=1.0, user-scalable=no">

    <title>Eka Automation Platform - Presentation</title>

    <link rel="stylesheet" href="https://cdnjs.cloudflare.com/ajax/libs/reveal.js/5.1.0/reset.min.css">
    <link rel="stylesheet" href="https://cdnjs.cloudflare.com/ajax/libs/reveal.js/5.1.0/reveal.min.css">
    <link rel="stylesheet" href="https://cdnjs.cloudflare.com/ajax/libs/reveal.js/5.1.0/theme/black.min.css">
    <link href="https://fonts.googleapis.com/icon?family=Material+Icons+Round" rel="stylesheet">
    <link href="https://fonts.googleapis.com/css2?family=Montserrat:wght@400;500;600;700;800&family=Poppins:wght@400;500;600;700&display=swap" rel="stylesheet">

    <style>
        body {
            background-color: #05080f !important;
            font-family: 'Poppins', sans-serif;
            color: #ffffff;
            overflow: hidden;
        }
        .reveal {
            font-family: 'Poppins', sans-serif;
            color: #ffffff;
        }
        .reveal h1, .reveal h2, .reveal h3, .reveal h4 {
            font-family: 'Montserrat', sans-serif;
            text-transform: none;
            color: #00f0ff;
            font-weight: 700;
            text-shadow: 0 0 10px rgba(0, 240, 255, 0.2);
            letter-spacing: -0.01em;
        }
        .reveal h1 {
            font-size: 2.8em;
            line-height: 1.1;
            font-weight: 800;
            margin-bottom: 12px;
            color: #ffffff;
        }
        .reveal h2 {
            color: #00f0ff;
            font-size: 1.7em;
            margin-bottom: 20px;
            border-bottom: 1px solid #1a2f4c;
            padding-bottom: 10px;
            text-align: left;
        }
        .slide-card {
            background: #0e1624;
            border: 1px solid #1a2f4c;
            border-radius: 12px;
            padding: 24px 30px;
            box-shadow: 0 12px 30px rgba(0, 0, 0, 0.5);
            text-align: left;
            margin-top: 10px;
            position: relative;
        }
        .slide-card::before {
            content: '';
            position: absolute;
            top: 0;
            left: 0;
            width: 100%;
            height: 3px;
            background: linear-gradient(90deg, #00f0ff, #007aff);
            border-radius: 12px 12px 0 0;
        }
        .card-title {
            color: #00f0ff;
            font-family: 'Montserrat', sans-serif;
            font-size: 1.15em;
            font-weight: 700;
            margin-bottom: 18px;
            display: flex;
            align-items: center;
            gap: 12px;
            border-bottom: 1px dashed rgba(26, 47, 76, 0.6);
            padding-bottom: 8px;
            letter-spacing: 0.5px;
            text-transform: uppercase;
        }
        .reveal ul {
            font-size: 0.65em;
            line-height: 1.7;
            list-style-type: none;
            margin: 0;
            padding: 0;
        }
        .reveal li {
            margin-bottom: 12px;
            position: relative;
            padding-left: 26px;
            color: #d1d9e6;
        }
        .reveal li::before {
            content: "✦";
            color: #00f0ff;
            text-shadow: 0 0 8px rgba(0, 240, 255, 0.8);
            position: absolute;
            left: 0;
            top: -1px;
        }
        .two-column {
            display: flex;
            flex-direction: row;
            align-items: center;
            gap: 30px;
            width: 100%;
        }
        .col-left {
            flex: 1;
            display: flex;
            flex-direction: column;
            justify-content: flex-start;
        }
        .col-right {
            flex: 1.2;
            position: relative;
            display: flex;
            align-items: center;
            justify-content: center;
        }
        .gallery-frame {
            background: #0e1624;
            border: 2px solid #00f0ff;
            padding: 8px;
            box-shadow: 0 15px 40px rgba(0, 240, 255, 0.15);
            border-radius: 8px;
            width: 100%;
            display: block;
        }
        .image-frame {
            width: 100%;
            border-radius: 4px;
            border: 1px solid #1a2f4c;
            display: block;
        }
        .arch-container {
            display: flex;
            flex-direction: column;
            align-items: center;
            gap: 15px;
            width: 100%;
            margin-top: 10px;
        }
        .arch-engine {
            background: linear-gradient(135deg, #00f0ff, #007aff);
            color: #05080f;
            padding: 14px 28px;
            border-radius: 8px;
            font-weight: 700;
            font-family: 'Montserrat', sans-serif;
            font-size: 0.8em;
            box-shadow: 0 0 20px rgba(0, 240, 255, 0.4);
            border: 1px solid #ffffff;
            text-align: center;
            letter-spacing: 0.5px;
        }
        .arch-grid-row1 {
            display: grid;
            grid-template-columns: repeat(4, 1fr);
            gap: 12px;
            width: 100%;
        }
        .arch-grid-row2 {
            display: flex;
            justify-content: center;
            gap: 15px;
            width: 100%;
        }
        .arch-pillar {
            background: #0e1624;
            border: 1px solid #1a2f4c;
            border-radius: 8px;
            padding: 12px;
            text-align: center;
            box-shadow: 0 4px 15px rgba(0, 0, 0, 0.3);
            transition: all 0.4s ease;
        }
        .arch-pillar:hover {
            border-color: #00f0ff;
            transform: translateY(-2px);
            box-shadow: 0 10px 25px rgba(0, 240, 255, 0.15);
        }
        .arch-pillar-title {
            color: #00f0ff;
            font-weight: 700;
            font-size: 0.6em;
            margin-bottom: 6px;
            font-family: 'Montserrat', sans-serif;
            letter-spacing: 0.5px;
        }
        .arch-pillar-desc {
            color: #a0aec0;
            font-size: 0.45em;
            line-height: 1.4;
        }
        .arch-line {
            width: 2px;
            height: 20px;
            background: #1a2f4c;
        }
        .download-btn {
            display: inline-flex;
            align-items: center;
            gap: 10px;
            background: linear-gradient(135deg, #00f0ff, #007aff);
            color: #05080f !important;
            padding: 12px 24px;
            border-radius: 6px;
            font-size: 0.52em !important;
            font-family: 'Montserrat', sans-serif;
            font-weight: 700;
            text-decoration: none;
            box-shadow: 0 6px 20px rgba(0, 240, 255, 0.25);
            transition: all 0.2s ease;
            margin-top: 25px;
        }
        .download-btn:hover {
            transform: translateY(-1px);
            box-shadow: 0 10px 25px rgba(0, 240, 255, 0.4);
        }
        .workflow-step {
            display: flex;
            align-items: center;
            gap: 16px;
            margin-bottom: 12px;
            text-align: left;
        }
        .workflow-number {
            width: 32px;
            height: 32px;
            background: #00f0ff;
            color: #05080f;
            border-radius: 50%;
            display: flex;
            align-items: center;
            justify-content: center;
            font-weight: 700;
            font-family: 'Montserrat', sans-serif;
            font-size: 0.7em;
            box-shadow: 0 0 10px rgba(0, 240, 255, 0.4);
            flex-shrink: 0;
        }
        .workflow-text {
            font-size: 0.6em;
            line-height: 1.4;
        }
        .workflow-text strong {
            color: #00f0ff;
            display: block;
        }
    </style>
</head>
<body>
    <div class="reveal">
        <div class="slides">

            <!-- Slide 1: Title -->
            <section style="text-align: left;">
                <div style="display: flex; gap: 40px; align-items: stretch; margin-left: 20px;">
                    <div style="width: 6px; background: linear-gradient(#00f0ff, #007aff); border-radius: 3px; box-shadow: 0 0 15px rgba(0, 240, 255, 0.4);"></div>
                    <div>
                        <h1 style="margin: 0; padding: 0; font-size: 2.7em; font-weight: 800;">EKA AUTOMATION PLATFORM</h1>
                        <p style="color: #00f0ff; font-size: 0.85em; font-family: 'Montserrat', sans-serif; font-weight: 600; margin: 15px 0 25px 0; text-shadow: 0 0 10px rgba(0,240,255,0.2);">
                            Enterprise-Grade One-Click Network Testbeds & Simulation Control Plane
                        </p>
                        <p style="color: #8899b8; font-size: 0.45em; letter-spacing: 1.5px; text-transform: uppercase; margin-bottom: 25px;">System Operations & Architecture Review</p>
                        
                        <a href="Eka_Automation_Presentation.pptx" download class="download-btn">
                            <span class="material-icons-round" style="font-size: 20px; vertical-align: middle;">download</span> 
                            Download PowerPoint (.pptx)
                        </a>
                    </div>
                </div>
            </section>

            <!-- Slide 2: Agenda -->
            <section>
                <h2>Presentation Agenda</h2>
                <div class="two-column">
                    <div class="col-left">
                        <div class="slide-card" style="height: 100%;">
                            <div class="card-title"><span class="material-icons-round" style="color: #00f0ff;">menu_book</span> Operational Core Elements</div>
                            <ul>
                                <li><strong>What is Eka?</strong> (Origin, Vision, and Control Plane)</li>
                                <li><strong>Problem Statement</strong> (Laboratory Staging Roadblocks)</li>
                                <li><strong>Solution Overview</strong> (Unified Simulation & Imaging Engine)</li>
                                <li><strong>Modular Architecture</strong> (The 7 Core Subsystem Tabs)</li>
                                <li><strong>Workflow Pipeline</strong> (Process Step-by-Step Walkthrough)</li>
                                <li><strong>Benefits & Demo</strong> (Time Savings & Interactive Screenshots)</li>
                            </ul>
                        </div>
                    </div>
                    <div class="col-right">
                        <div class="gallery-frame">
                            <img src="screenshot_dashboard.png" class="image-frame" alt="Eka Dashboard">
                        </div>
                    </div>
                </div>
            </section>

            <!-- Slide 3: Problem Statement -->
            <section style="text-align: left;">
                <h2>Problem Statement: Laboratory Friction</h2>
                <div class="slide-card">
                    <div class="card-title" style="color: #ff3838; text-shadow: 0 0 10px rgba(255,56,56,0.2);">
                        <span class="material-icons-round">error_outline</span> Staging & QA Complexities
                    </div>
                    <ul>
                        <li><strong>Manual Staging Delays:</strong> Traditional virtual machines and physical switches staging take hours of complex command-line entry.</li>
                        <li><strong>Configuration Drift:</strong> Copying firmware versions manually and accessing hypervisors directly leads to error-prone deployments.</li>
                        <li><strong>Tooling Fragmentation:</strong> Test developers rely on multiple programs (PuTTY, SCP file copy clients, terminal multiplexers).</li>
                        <li><strong>Resource Collision:</strong> Multiple QA engineers attempting to connect or flash identical switches at the same time.</li>
                    </ul>
                </div>
            </section>

            <!-- Slide 4: Solution Overview -->
            <section style="text-align: left;">
                <h2>Solution: Unified Control Plane</h2>
                <div class="slide-card">
                    <div class="card-title"><span class="material-icons-round">check_circle</span> Automated Orchestration Plane</div>
                    <ul>
                        <li><strong>Unified Control:</strong> Eka pools hypervisors, terminal SSH clients, and switch re-imaging loops into a central browser view.</li>
                        <li><strong>Automated Provisioning:</strong> Asynchronous SCP copies and flashes firmware packages to targets automatically.</li>
                        <li><strong>Direct Cable Topologies:</strong> Draw virtual Ethernet inter-interface cable connections directly on an execution canvas.</li>
                        <li><strong>Heartbeat Locking:</strong> Background resource locking mechanisms block device connection overlap.</li>
                    </ul>
                </div>
            </section>

            <!-- Slide 5: Modular System Architecture -->
            <section>
                <h2>Modular System Architecture</h2>
                <div class="arch-container">
                    <div class="arch-engine">
                        EKA CONTROL PLANE ORCHESTRATOR
                    </div>
                    <div class="arch-line"></div>
                    <div class="arch-grid-row1">
                        <div class="arch-pillar">
                            <div class="arch-pillar-title">DASHBOARD</div>
                            <div class="arch-pillar-desc">Real-time stats visualizer.</div>
                        </div>
                        <div class="arch-pillar">
                            <div class="arch-pillar-title">DEVICES</div>
                            <div class="arch-pillar-desc">Switch and VM register list.</div>
                        </div>
                        <div class="arch-pillar">
                            <div class="arch-pillar-title">EXECUTION</div>
                            <div class="arch-pillar-desc">Interactive topology canvas.</div>
                        </div>
                        <div class="arch-pillar">
                            <div class="arch-pillar-title">VS MANAGER</div>
                            <div class="arch-pillar-desc">Hypervisor VM coordinator.</div>
                        </div>
                    </div>
                    <div class="arch-line"></div>
                    <div class="arch-grid-row2">
                        <div class="arch-pillar" style="width: 28%;">
                            <div class="arch-pillar-title">AUDIT LOGS</div>
                            <div class="arch-pillar-desc">Trace files repository database.</div>
                        </div>
                        <div class="arch-pillar" style="width: 28%;">
                            <div class="arch-pillar-title">SSH TERMINAL</div>
                            <div class="arch-pillar-desc">Web browser SSH sessions.</div>
                        </div>
                        <div class="arch-pillar" style="width: 28%;">
                            <div class="arch-pillar-title">HARDWARE LOAD</div>
                            <div class="arch-pillar-desc">Physical switch firmware installer.</div>
                        </div>
                    </div>
                </div>
            </section>

            <!-- Slide 6: Key Features -->
            <section style="text-align: left;">
                <h2>Eka Core Feature Highlights</h2>
                <div class="slide-card">
                    <div class="card-title"><span class="material-icons-round">star_outline</span> Powering Modern QA Operations</div>
                    <ul>
                        <li><strong>Drag-and-Drop Topology:</strong> Draw cabling paths between switches on a GNS3-style interactive canvas.</li>
                        <li><strong>ONIE Loop Detection:</strong> Identifies switches stuck in boot loops and flash firmware asynchronously.</li>
                        <li><strong>Persistent Web SSH:</strong> Multiplex SSH console commands without third-party tools; persists commands in tmux sessions.</li>
                        <li><strong>Asynchronous SCP Copying:</strong> Update VM and physical switches concurrently via multi-threaded download pools.</li>
                        <li><strong>Logs Archival database:</strong> Search completed runs, logs output, and execution history.</li>
                    </ul>
                </div>
            </section>

            <!-- Slide 7: Automation Pipeline -->
            <section>
                <h2>Workflow & Automation Pipeline</h2>
                <div class="slide-card">
                    <div class="card-title"><span class="material-icons-round">alt_route</span> Orchestration Process Flow</div>
                    <div style="margin-top: 15px;">
                        <div class="workflow-step">
                            <div class="workflow-number">1</div>
                            <div class="workflow-text">
                                <strong>Register Switches</strong>
                                Onboard physical switch and hypervisor IP credentials in the Devices register.
                            </div>
                        </div>
                        <div class="workflow-step">
                            <div class="workflow-number">2</div>
                            <div class="workflow-text">
                                <strong>Flash Builds</strong>
                                Initiate SCP copy transfers to flash switch boot loaders or virtual system disks.
                            </div>
                        </div>
                        <div class="workflow-step">
                            <div class="workflow-number">3</div>
                            <div class="workflow-text">
                                <strong>Draw Topology</strong>
                                Connect interface ports on the Canvas to construct the target regression topography.
                            </div>
                        </div>
                        <div class="workflow-step">
                            <div class="workflow-number">4</div>
                            <div class="workflow-text">
                                <strong>Run SpyTest</strong>
                                Dispatch selected Python script batches and watch execution logs stream.
                            </div>
                        </div>
                        <div class="workflow-step">
                            <div class="workflow-number">5</div>
                            <div class="workflow-text">
                                <strong>Trace Audits</strong>
                                Store all run outputs and standard output text in the Logs query database.
                            </div>
                        </div>
                    </div>
                </div>
            </section>

            <!-- Slide 8: Platform Benefits -->
            <section style="text-align: left;">
                <h2>Platform Benefits: Dev & QA Teams</h2>
                <div class="slide-card">
                    <div class="card-title"><span class="material-icons-round">trending_up</span> Engineering Advantages</div>
                    <ul>
                        <li><strong>Drastic Time Savings:</strong> Automatic ONIE loops and hypervisor flash workflows save up to 45 minutes per setup cycle.</li>
                        <li><strong>Simplified Access:</strong> Access switch CLI terminal sessions from any browser; zero VPN or client programs (PuTTY) required.</li>
                        <li><strong>Zero Configuration Drift:</strong> Flashes clean builds and clears residual config files before dispatching runs.</li>
                        <li><strong>Efficient Resource Pools:</strong> Visual device reservation indicators prevent hardware overlap and collision.</li>
                    </ul>
                </div>
            </section>

            <!-- Slide 9: Scalability & Performance -->
            <section style="text-align: left;">
                <h2>Scalability & Performance Metrics</h2>
                <div class="slide-card">
                    <div class="card-title"><span class="material-icons-round">speed</span> Engineered for Velocity</div>
                    <ul>
                        <li><strong>Multi-Threaded Copy:</strong> Simultaneous SCP downloads flash up to 8 switches concurrently without speed drops.</li>
                        <li><strong>FastAPI Async Engine:</strong> Lightweight heartbeats fetch status parameters with under 2% controller overhead.</li>
                        <li><strong>Active Paramiko Pools:</strong> Re-uses backend SSH pools to connect and send CLI commands in under 100ms.</li>
                        <li><strong>Distributed Docker Scaling:</strong> Server components can be scaled out in containers to manage hundreds of active nodes.</li>
                    </ul>
                </div>
            </section>

            <!-- Slides 10-16: Live Tab Walkthroughs -->
            <!-- Slide 10: Dashboard -->
            <section>
                <h2>1. Dashboard Tab</h2>
                <div class="two-column">
                    <div class="col-left">
                        <div class="slide-card" style="height: 100%;">
                            <div class="card-title"><span class="material-icons-round">dashboard</span> Control Panel Overview</div>
                            <ul>
                                <li>Monitor live connection states of all registered devices in real time.</li>
                                <li>Display visual counters for total devices, online nodes, and script queues.</li>
                                <li>Review run outcomes via recent execution telemetry charts.</li>
                            </ul>
                        </div>
                    </div>
                    <div class="col-right">
                        <div class="gallery-frame">
                            <img src="screenshot_dashboard.png" class="image-frame" alt="Eka Dashboard">
                        </div>
                    </div>
                </div>
            </section>

            <!-- Slide 11: Devices -->
            <section>
                <h2>2. Devices Tab</h2>
                <div class="two-column">
                    <div class="col-left">
                        <div class="slide-card" style="height: 100%;">
                            <div class="card-title"><span class="material-icons-round">dns</span> Inventory Management</div>
                            <ul>
                                <li>Onboard physical switches and virtual machines with IP credentials.</li>
                                <li>Modify network parameter mappings (Ports, SSH credentials, paths).</li>
                                <li>Perform background ping heartbeats to determine device active states.</li>
                            </ul>
                        </div>
                    </div>
                    <div class="col-right">
                        <div class="gallery-frame">
                            <img src="screenshot_devices.png" class="image-frame" alt="Eka Devices">
                        </div>
                    </div>
                </div>
            </section>

            <!-- Slide 12: Execution -->
            <section>
                <h2>3. Execution Tab</h2>
                <div class="two-column">
                    <div class="col-left">
                        <div class="slide-card" style="height: 100%;">
                            <div class="card-title"><span class="material-icons-round">play_circle_outline</span> Interactive Topology</div>
                            <ul>
                                <li>Draw cabling connections between virtual device interfaces on a canvas.</li>
                                <li>Browse hierarchical test script folders and select test groups.</li>
                                <li>Trigger automated SpyTest executions with a single click.</li>
                            </ul>
                        </div>
                    </div>
                    <div class="col-right">
                        <div class="gallery-frame">
                            <img src="screenshot_execute.png" class="image-frame" alt="Eka Execution">
                        </div>
                    </div>
                </div>
            </section>

            <!-- Slide 13: VS Manager -->
            <section>
                <h2>4. VS Manager Tab</h2>
                <div class="two-column">
                    <div class="col-left">
                        <div class="slide-card" style="height: 100%;">
                            <div class="card-title"><span class="material-icons-round">settings_system_daydream</span> Hypervisor Coordination</div>
                            <ul>
                                <li>Query and display VM states from remote hypervisors.</li>
                                <li>Power on/off individual or groups of virtual switches instantly.</li>
                                <li>Fetch and copy VM build images from network storage via SCP.</li>
                            </ul>
                        </div>
                    </div>
                    <div class="col-right">
                        <div class="gallery-frame">
                            <img src="screenshot_vs.png" class="image-frame" alt="Eka VS Manager">
                        </div>
                    </div>
                </div>
            </section>

            <!-- Slide 14: Logs -->
            <section>
                <h2>5. Logs Tab</h2>
                <div class="two-column">
                    <div class="col-left">
                        <div class="slide-card" style="height: 100%;">
                            <div class="card-title"><span class="material-icons-round">history</span> Trace Logs Database</div>
                            <ul>
                                <li>Access database audit logs for all completed and failed runs.</li>
                                <li>Expand execution rows to inspect live terminal stdio streams.</li>
                                <li>Filter records by date, script name, or test engineer.</li>
                            </ul>
                        </div>
                    </div>
                    <div class="col-right">
                        <div class="gallery-frame">
                            <img src="screenshot_logs.png" class="image-frame" alt="Eka Logs">
                        </div>
                    </div>
                </div>
            </section>

            <!-- Slide 15: Terminal -->
            <section>
                <h2>6. Terminal Tab</h2>
                <div class="two-column">
                    <div class="col-left">
                        <div class="slide-card" style="height: 100%;">
                            <div class="card-title"><span class="material-icons-round">terminal</span> Integrated SSH Consoles</div>
                            <ul>
                                <li>Open instant SSH terminal sessions to any registered device inside the UI.</li>
                                <li>Execute standard CLI commands without leaving the web browser.</li>
                                <li>Preserve console states using background tmux session persistence.</li>
                            </ul>
                        </div>
                    </div>
                    <div class="col-right">
                        <div class="gallery-frame">
                            <img src="screenshot_terminal.png" class="image-frame" alt="Eka Terminal">
                        </div>
                    </div>
                </div>
            </section>

            <!-- Slide 16: Hardware Load -->
            <section>
                <h2>7. Hardware Load Tab</h2>
                <div class="two-column">
                    <div class="col-left">
                        <div class="slide-card" style="height: 100%;">
                            <div class="card-title"><span class="material-icons-round">settings_input_hdmi</span> ONIE Firmware Flasher</div>
                            <ul>
                                <li>Update switch builds automatically via a single-button installer.</li>
                                <li>Auto-detect switches in ONIE boot loops and flash firmware packages.</li>
                                <li>View real-time console boot sequence logs during re-imaging.</li>
                            </ul>
                        </div>
                    </div>
                    <div class="col-right">
                        <div class="gallery-frame">
                            <img src="screenshot_hardware_load.png" class="image-frame" alt="Eka Hardware Load">
                        </div>
                    </div>
                </div>
            </section>

            <!-- Slide 17: Roadmap -->
            <section style="text-align: left;">
                <h2>Conclusion & Future Roadmap</h2>
                <div class="slide-card">
                    <div class="card-title"><span class="material-icons-round">map</span> Expanding Automation Capabilities</div>
                    <ul>
                        <li><strong>OIDC SSO Auth:</strong> Delegate authorization to corporate platforms (OnePalc portal redirect).</li>
                        <li><strong>Multi-Hypervisor Canvas:</strong> Connect virtual networks spanning different remote ESXi/KVM hosts.</li>
                        <li><strong>Machine Learning Diagnostics:</strong> Automatic classification of test run logs to pinpoint root failure causes.</li>
                        <li><strong>Visual Device Health:</strong> Real-time temperature, memory, and CPU metric charts on the dashboard.</li>
                    </ul>
                </div>
            </section>

            <!-- End Slide -->
            <section style="text-align: center;">
                <h1 style="font-size: 3em; color: #00f0ff;">THANK YOU.</h1>
                <p style="color: #8899b8; font-size: 1.1em; font-family: 'Montserrat', sans-serif; font-weight: 600;">Innovative Test Automation Orchestration</p>
                <div style="font-size: 0.38em; color: #5e685f; margin-top: 30px; font-style: italic;">
                    Press ESC to see slide overview. Use arrow keys to navigate.
                </div>
            </section>

        </div>
    </div>

    <script src="https://cdnjs.cloudflare.com/ajax/libs/reveal.js/5.1.0/reveal.min.js"></script>
    <script>
        Reveal.initialize({
            hash: true,
            slideNumber: 'c/t',
            progress: true,
            center: true,
            width: 1280,
            height: 720,
            transition: 'slide',
            backgroundTransition: 'fade',
            controls: true,
            keyboard: true
        });
    </script>
</body>
</html>
"""
    os.makedirs("static", exist_ok=True)

    with open("Eka_Automation_Web_Presentation.html", "w", encoding="utf-8") as f:
        f.write(html_content)
    print("[OK] Saved HTML to Eka_Automation_Web_Presentation.html")

    # Write to static folder
    with open("static/Eka_Automation_Web_Presentation.html", "w", encoding="utf-8") as f:
        f.write(html_content)
    print("[OK] Saved HTML to static/Eka_Automation_Web_Presentation.html")

    # Copy screenshots to static/
    images_to_copy = [
        "screenshot_dashboard.png", "screenshot_devices.png", "screenshot_execute.png", 
        "screenshot_vs.png", "screenshot_logs.png", "screenshot_terminal.png", "screenshot_hardware_load.png"
    ]
    for img in images_to_copy:
        if os.path.exists(img):
            shutil.copy(img, os.path.join("static", img))
            print(f"[OK] Copied {img} to static/")
            
    # Copy PPTX to static/
    if os.path.exists("Eka_Automation_Presentation.pptx"):
        shutil.copy("Eka_Automation_Presentation.pptx", os.path.join("static", "Eka_Automation_Presentation.pptx"))
        print("[OK] Copied Eka_Automation_Presentation.pptx to static/")

if __name__ == "__main__":
    print("Compiling Dark Futuristic PPTX Presentation...")
    compile_pptx()
    print("Compiling Dark Futuristic HTML Presentation...")
    compile_html()
    print("Done generating presentation files successfully!")
